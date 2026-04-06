"""
Niu File Parser MCP Server

Parses various document formats (PDF, Word, PPT, Excel, Markdown, HTML)
and returns structured content for knowledge graph ingestion.
"""

import asyncio
import json
from pathlib import Path
from typing import Any

from loguru import logger
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import TextContent, Tool

# Initialize MCP server
server = Server("niu-file-parser")


def parse_pdf(file_path: str) -> dict[str, Any]:
    """Parse PDF file and extract text content."""
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages.append({"page": i + 1, "content": text.strip()})

    return {
        "file_type": "pdf",
        "total_pages": len(reader.pages),
        "pages": pages,
    }


def parse_docx(file_path: str) -> dict[str, Any]:
    """Parse Word document and extract text content."""
    from docx import Document

    doc = Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())

    return {
        "file_type": "docx",
        "paragraphs": paragraphs,
        "total_paragraphs": len(paragraphs),
    }


def parse_pptx(file_path: str) -> dict[str, Any]:
    """Parse PowerPoint file and extract text content."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            # Use getattr to safely access text attribute
            text = getattr(shape, "text", "")
            if text and str(text).strip():
                texts.append(str(text).strip())
        slides.append({"slide": i + 1, "content": "\n".join(texts)})

    return {
        "file_type": "pptx",
        "total_slides": len(slides),
        "slides": slides,
    }


def parse_xlsx(file_path: str) -> dict[str, Any]:
    """Parse Excel file and extract text content."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_data = [str(cell) if cell is not None else "" for cell in row]
            if any(row_data):  # Skip empty rows
                rows.append(row_data)
        sheets.append({"name": sheet_name, "rows": rows})

    return {
        "file_type": "xlsx",
        "sheets": sheets,
    }


def parse_markdown(file_path: str) -> dict[str, Any]:
    """Parse Markdown file."""
    content = Path(file_path).read_text(encoding="utf-8")
    return {
        "file_type": "markdown",
        "content": content,
    }


def parse_html(file_path: str) -> dict[str, Any]:
    """Parse HTML file and extract text content."""
    from bs4 import BeautifulSoup

    content = Path(file_path).read_text(encoding="utf-8")
    soup = BeautifulSoup(content, "lxml")

    # Remove script and style elements
    for element in soup(["script", "style", "nav", "footer", "header"]):
        element.decompose()

    text = soup.get_text(separator="\n", strip=True)

    return {
        "file_type": "html",
        "title": soup.title.string if soup.title else "",
        "content": text,
    }


def parse_text(file_path: str) -> dict[str, Any]:
    """Parse plain text file."""
    content = Path(file_path).read_text(encoding="utf-8")
    return {
        "file_type": "text",
        "content": content,
    }


@server.list_tools()
async def list_tools() -> list[Tool]:
    """List available tools."""
    return [
        Tool(
            name="parse_file",
            description="Parse a document file (PDF, Word, PPT, Excel, Markdown, HTML, Text) and extract text content.",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Absolute path to the file to parse",
                    },
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="list_supported_formats",
            description="List all supported file formats for parsing.",
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict[str, Any]) -> list[TextContent]:
    """Handle tool calls."""
    if name == "list_supported_formats":
        formats = [
            {"extension": ".pdf", "description": "PDF documents"},
            {"extension": ".docx", "description": "Microsoft Word documents"},
            {"extension": ".pptx", "description": "Microsoft PowerPoint presentations"},
            {"extension": ".xlsx", "description": "Microsoft Excel spreadsheets"},
            {"extension": ".md", "description": "Markdown files"},
            {"extension": ".html", "description": "HTML files"},
            {"extension": ".txt", "description": "Plain text files"},
        ]
        return [TextContent(type="text", text=json.dumps(formats, indent=2))]

    if name == "parse_file":
        file_path = arguments.get("file_path", "")
        path = Path(file_path)

        if not path.exists():
            return [
                TextContent(type="text", text=f"Error: File not found: {file_path}")
            ]

        suffix = path.suffix.lower()

        try:
            if suffix == ".pdf":
                result = parse_pdf(file_path)
            elif suffix == ".docx":
                result = parse_docx(file_path)
            elif suffix == ".pptx":
                result = parse_pptx(file_path)
            elif suffix == ".xlsx":
                result = parse_xlsx(file_path)
            elif suffix == ".md":
                result = parse_markdown(file_path)
            elif suffix == ".html":
                result = parse_html(file_path)
            elif suffix == ".txt":
                result = parse_text(file_path)
            else:
                return [
                    TextContent(
                        type="text",
                        text=f"Error: Unsupported file format: {suffix}. Use list_supported_formats to see supported formats.",
                    )
                ]

            result["file_path"] = file_path
            result["file_name"] = path.name
            return [
                TextContent(
                    type="text", text=json.dumps(result, indent=2, ensure_ascii=False)
                )
            ]

        except Exception as e:
            logger.exception(f"Error parsing file: {e}")
            return [TextContent(type="text", text=f"Error parsing file: {e}")]

    return [TextContent(type="text", text=f"Unknown tool: {name}")]


async def run_server():
    """Run the MCP server."""
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream, write_stream, server.create_initialization_options()
        )


def main():
    """Main entry point."""
    asyncio.run(run_server())


if __name__ == "__main__":
    main()
