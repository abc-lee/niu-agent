"""
Niu File & Photo MCP Server

Provides tools for file and photo management.
Supports document ingestion, photo processing with face recognition.
"""

import asyncio
import hashlib
import json
import os
import shutil
import sqlite3
import threading
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

import numpy as np
from loguru import logger
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import TextContent, Tool

# ============== Tool Schemas ==============

TOOL_SCHEMAS = {
    "ingest": {
        "name": "ingest",
        "description": """有状态统一入库工具 — 支持目录逐文件交互式入库

参数:
- path: 必填，文件路径或目录路径
- mode: copy（复制）| move（移动）| reference（引用），默认 copy
- category: 分类目录（文档需要分类时传入，照片不需要）
- action: start | interact | abort，默认空字符串

三阶段交互模式（目录入库）:
1. 初始化: ingest(path="E:/照片", action="start", mode="copy")
   → 扫描目录，创建会话
   → 如果传入 category：所有文档使用同一分类，自动循环处理到完成
   → 如果未传 category：自动处理所有照片，遇到文档时返回 need_category
2. 中间态交互:
   - 回答分类（need_category后）: ingest(path="E:/照片", category="技术文档")
   → 仅处理当前文档，然后继续自动处理照片，遇到下一个未分类文档再次返回 need_category
3. 中止: ingest(path="E:/照片", action="abort")

三种自动循环场景:
- 初始化带 category → 自动循环到完成（所有文档用同一分类）
- 纯照片目录（无文档）→ 自动循环到完成
- 子Agent回答分类后 → 只处理当前文档，不会自动跳过后续未分类文档

单文件入库（path是文件时）: 无状态，直接入库，action参数无效

返回:
- 照片: {status, photo_id, detected_persons, abstract, exif}
- 文档(need_category): {status: need_category, preview, available_categories}
- 文档(success): {status: success, action, file_path, lightrag, lightrag_message}
- 目录(need_category): {status: need_category, total, current_file, preview, available_categories}
- 目录(success): {status: success, total, photos, documents, skipped, details}
- 目录(aborted): {status: aborted, message, processed_count}
- 错误: {status: error, message}""",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "文件路径或目录路径"},
                "mode": {
                    "type": "string",
                    "enum": ["copy", "move", "reference"],
                    "default": "copy",
                    "description": "文件操作模式",
                },
                "category": {"type": "string", "description": "文件分类目录。need_category状态时必须从available_categories中选择", "default": ""},
                "action": {
                    "type": "string",
                    "enum": ["", "start", "interact", "abort"],
                    "default": "",
                    "description": "会话动作：start=初始化会话，interact/空=继续交互，abort=中止会话",
                },
            },
            "required": ["path"],
        },
    },
    "ingest_document": {
        "name": "ingest_document",
        "description": """文档入库工具 — 文件搬运 + 提交LightRAG异步处理

参数:
- file_path: 必填，文档文件路径
- category: 分类目录，不传则返回内容预览供判断分类
- mode: copy（复制）| move（移动）| reference（引用），默认 copy

不传 category 时返回 need_category 状态+内容预览，判断分类后再次调用 ingest 或 ingest_document 传入 category。

返回:
- status: need_category | success | error
- action: created | versioned | renamed | referenced | skipped
- file_path: 存储路径
- lightrag: inserted | unsupported | skipped | error
- lightrag_message: unsupported/error 时的原因说明""",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "文档文件路径"},
                "category": {"type": "string", "description": "文件分类目录（如：工作文档、个人资料、财务报告等）。不传则返回文件内容预览供你判断分类", "default": ""},
                "mode": {
                    "type": "string",
                    "enum": ["copy", "move", "reference"],
                    "default": "copy",
                    "description": "文件操作模式",
                },
            },
            "required": ["file_path"],
        },
    },
    "name_person": {
        "name": "name_person",
        "description": """为人物命名

参数:
- person_id: 必填，人物ID
- name: 必填，新名称

返回:
- status: success | error
- person_id: 人物ID
- name: 新名称
- auto_label: 自动标签（如"未命名人物_1"）""",
        "input_schema": {
            "type": "object",
            "properties": {
                "person_id": {"type": "string", "description": "人物ID"},
                "name": {"type": "string", "description": "新名称"},
            },
            "required": ["person_id", "name"],
        },
    },
    "merge_persons": {
        "name": "merge_persons",
        "description": """合并两个人物

参数:
- person_a_id: 必填，保留的人物ID
- person_b_id: 必填，要合并到A的人物ID

返回:
- status: success | error
- merged_into: 合并到的人物ID
- name: 保留的名称
- photo_count: 合并后的照片数量
- deleted_person_id: 被删除的人物ID

说明:
- 保留 person_a 的名称
- 合并所有人脸向量，重新计算中心
- 更新所有照片关联
- 学习机制：如果相似度低于阈值，自动调整阈值""",
        "input_schema": {
            "type": "object",
            "properties": {
                "person_a_id": {"type": "string", "description": "保留的人物ID"},
                "person_b_id": {"type": "string", "description": "要合并的人物ID"},
            },
            "required": ["person_a_id", "person_b_id"],
        },
    },
    "search_persons": {
        "name": "search_persons",
        "description": """搜索人物（按名字模糊匹配）

参数:
- query: 搜索词（人名，子串匹配）
- limit: 返回数量（默认10）

返回:
- 匹配的人物列表""",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "搜索词（人名）"},
                "limit": {
                    "type": "integer",
                    "description": "返回数量",
                    "default": 10,
                },
            },
            "required": ["query"],
        },
    },
    "get_unnamed_persons": {
        "name": "get_unnamed_persons",
        "description": """获取所有未命名人物

返回:
- 未命名人物列表，按出现次数排序
- 包含：id, auto_label, photo_count, has_valid_photos, photos: [{file_path, boxed_path}]
- boxed_path 是带人脸红框的图片路径，前端用 Markdown 图片语法 ![person_id|name](path) 显示
- has_valid_photos=false 表示该人物的照片文件已不存在""",
        "input_schema": {
            "type": "object",
            "properties": {},
        },
    },
    "delete_person": {
        "name": "delete_person",
        "description": """删除人物及其所有关联数据

警告：这会删除人物图谱中的节点，请谨慎使用。
只有在用户明确要求删除时才调用。

参数:
- person_id: 要删除的人物ID

返回:
- status: success | error
- message: 结果说明""",
        "input_schema": {
            "type": "object",
            "properties": {
                "person_id": {"type": "string", "description": "人物ID"},
            },
            "required": ["person_id"],
        },
    },
    "cleanup_deleted_photos": {
        "name": "cleanup_deleted_photos",
        "description": """清理已删除照片的数据库记录

在删除照片文件/目录后调用，清理数据库中的孤儿记录。
扫描 photos 表，删除文件不存在的记录及其关联的 faces 记录。

返回:
- deleted_photos: 删除的照片记录数
- deleted_faces: 删除的人脸记录数

使用场景：
- 用户删除了照片目录后
- 清理数据库中的残留记录""",
        "input_schema": {
            "type": "object",
            "properties": {},
        },
    },
    "get_person_photos": {
        "name": "get_person_photos",
        "description": """获取某人物的多张照片（用于"换一张"场景）

当用户说"看不清"、"换一张"时调用此工具。

参数:
- person_id: 人物ID
- limit: 最多返回几张照片（默认5）

返回:
- person_id, person_name
- photos: [{file_path, boxed_path, taken_at}, ...]

boxed_path 是带人脸红框的图片路径，前端用 Markdown 图片语法 ![person_id|name](path) 显示。""",
        "input_schema": {
            "type": "object",
            "properties": {
                "person_id": {"type": "string", "description": "人物ID"},
                "limit": {
                    "type": "integer",
                    "description": "最多返回几张（默认5）",
                },
            },
            "required": ["person_id"],
        },
    },
    "unload_face_model": {
        "name": "unload_face_model",
        "description": """卸载人脸识别模型，释放内存（约 326MB）

在长时间空闲时调用此工具释放内存。
通常由系统在 SLEEP 状态时自动调用。

返回:
- status: success
- message: 卸载结果""",
        "input_schema": {
            "type": "object",
            "properties": {},
        },
    },
}


def get_tool_schemas() -> list[dict]:
    """返回所有工具的 schema 列表（用于 MCP Loader 注册）"""
    return list(TOOL_SCHEMAS.values())


# ============== Database ==============


def get_db_path() -> Path:
    """Get photo database path (3-level priority, consistent with get_vector_db_path).

    Priority:
    1. NIU_DB_PATH env var (explicit override)
    2. WORKSPACE_PATH env var (set by Go launcher main.go)
    3. ~/.niu/memory.json workspace.path

    Auto-creates workspace directory if it doesn't exist.
    """
    # 1. NIU_DB_PATH — explicit override
    if "NIU_DB_PATH" in os.environ:
        p = Path(os.environ["NIU_DB_PATH"])
        if not p.is_absolute():
            raise ValueError(f"NIU_DB_PATH 必须是绝对路径: {p}")
        # 解析符号链接
        p = p.resolve()
        # 如果路径没有扩展名，假设是目录，添加 photos.db
        if p.suffix == '':
            p = p / "photos.db"
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
        except PermissionError as e:
            raise ValueError(f"无权限创建目录: {p.parent}。请检查权限。") from e
        except FileExistsError as e:
            raise ValueError(f"路径已存在但不是目录（可能是文件）: {p.parent}") from e
        return p
    # 2. WORKSPACE_PATH env var
    if "WORKSPACE_PATH" in os.environ:
        ws = Path(os.environ["WORKSPACE_PATH"])
        if not ws.is_absolute():
            raise ValueError(f"WORKSPACE_PATH 必须是绝对路径: {ws}")
        # 解析符号链接
        ws = ws.resolve()
        if ws.exists() and not ws.is_dir():
            raise ValueError(f"WORKSPACE_PATH 指向的不是目录: {ws}")
        try:
            ws.mkdir(parents=True, exist_ok=True)
        except PermissionError as e:
            raise ValueError(f"无权限创建工作区目录: {ws}。请检查权限。") from e
        except FileExistsError as e:
            raise ValueError(f"路径已存在但不是目录（可能是文件）: {ws}") from e
        return ws / "photos.db"
    # 3. ~/.niu/memory.json workspace.path
    memory_path = Path.home() / ".niu" / "memory.json"
    if memory_path.exists():
        try:
            with open(memory_path, "r", encoding="utf-8") as f:
                memory = json.load(f)
            workspace_path = memory.get("workspace", {}).get("path")
            if workspace_path and workspace_path.strip():
                ws = Path(workspace_path)
                if not ws.is_absolute():
                    raise ValueError(f"workspace.path 必须是绝对路径: {workspace_path}")
                # 解析符号链接
                ws = ws.resolve()
                if ws.exists() and not ws.is_dir():
                    raise ValueError(f"workspace.path 指向的不是目录: {workspace_path}")
                try:
                    ws.mkdir(parents=True, exist_ok=True)
                except PermissionError as e:
                    raise ValueError(f"无权限创建工作区目录: {ws}。请检查权限。") from e
                except FileExistsError as e:
                    raise ValueError(f"路径已存在但不是目录（可能是文件）: {ws}") from e
                return ws / "photos.db"
        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"无法从 {memory_path} 解析 JSON: {e}。请检查 memory.json 格式是否正确。") from e
    raise ValueError(
        f"无法确定照片库路径：~/.niu/memory.json 不存在或缺少 workspace.path 配置。"
        f"请在 ~/.niu/memory.json 中设置 workspace.path，或设置 WORKSPACE_PATH 环境变量。"
    )


_conn: sqlite3.Connection | None = None
_db_path_failed: bool = False
_db_write_lock = threading.Lock()

# ============== 有状态入库会话 ==============
_ingest_sessions: dict[str, dict] = {}


def get_connection() -> sqlite3.Connection:
    """Get or create database connection."""
    global _conn, _db_path_failed
    if _db_path_failed:
        raise RuntimeError("照片库路径解析失败，无法建立连接。请检查 memory.json 中 workspace.path 配置。")
    if _conn is None:
        try:
            db_path = get_db_path()
        except ValueError as e:
            _db_path_failed = True
            logger.error(f"照片库路径解析失败: {e}")
            raise RuntimeError(f"照片库路径解析失败: {e}") from e
        db_path.parent.mkdir(parents=True, exist_ok=True)
        _conn = sqlite3.connect(str(db_path), check_same_thread=False)
        _init_schema(_conn)
    return _conn


def _init_schema(conn: sqlite3.Connection) -> None:
    """Initialize database schema."""
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS persons (
            id TEXT PRIMARY KEY,
            name TEXT,
            auto_label TEXT,
            name_embedding BLOB,          -- 人物名向量（用于搜索）
            center_embedding BLOB,        -- 人脸中心向量
            threshold_adjustment REAL DEFAULT 0,
            photo_count INTEGER DEFAULT 0,
            first_seen TEXT,
            last_seen TEXT,
            created_at TEXT
        );
        
        CREATE TABLE IF NOT EXISTS photos (
            id TEXT PRIMARY KEY,
            file_path TEXT,
            taken_at TEXT,
            location TEXT,
            camera TEXT,
            abstract TEXT,
            ingested_at TEXT,
            kg_synced INTEGER DEFAULT 0
        );
        
        CREATE TABLE IF NOT EXISTS faces (
            id TEXT PRIMARY KEY,
            photo_id TEXT,
            person_id TEXT,
            embedding BLOB,
            bounding_box TEXT,
            confidence REAL,
            FOREIGN KEY (photo_id) REFERENCES photos(id),
            FOREIGN KEY (person_id) REFERENCES persons(id)
        );
        
        -- 同框关系表
        CREATE TABLE IF NOT EXISTS co_occurrences (
            person_a_id TEXT,
            person_b_id TEXT,
            count INTEGER DEFAULT 1,
            first_seen TEXT,
            last_seen TEXT,
            PRIMARY KEY (person_a_id, person_b_id),
            FOREIGN KEY (person_a_id) REFERENCES persons(id),
            FOREIGN KEY (person_b_id) REFERENCES persons(id)
        );
        
        CREATE INDEX IF NOT EXISTS idx_faces_photo ON faces(photo_id);
        CREATE INDEX IF NOT EXISTS idx_faces_person ON faces(person_id);
        CREATE INDEX IF NOT EXISTS idx_photos_taken ON photos(taken_at);
        CREATE INDEX IF NOT EXISTS idx_co_occurrences ON co_occurrences(person_a_id, person_b_id);
    """)
    # Migration: add kg_synced column if missing (existing DBs)
    try:
        conn.execute("ALTER TABLE photos ADD COLUMN kg_synced INTEGER DEFAULT 0")
    except sqlite3.OperationalError as e:
        if "duplicate column" not in str(e).lower():
            raise
    conn.commit()
    logger.info("Photo database schema initialized")


# ============== 共享向量服务 ==============
# 同进程架构：直接调用 niu_api.internal.embedding，不走 HTTP


def call_embedding_service(endpoint: str, data: dict) -> dict | None:
    """直接调用 niu_api.internal.embedding（同进程，无 HTTP 开销）。"""
    try:
        from niu_api.internal.embedding import encode, similarity

        if endpoint == "/encode":
            text = data.get("text", "")
            embedding = encode(text)
            return {"embedding": embedding}
        elif endpoint == "/similarity":
            text1 = data.get("text1", "")
            text2 = data.get("text2", "")
            sim = similarity(text1, text2)
            return {"similarity": sim}
    except Exception as e:
        logger.warning(f"[Embedding] 同进程调用失败: {e}")
        return None


# ============== 知识图谱同步 ==============


def _generate_stable_description(normalized_stem: str, abstract: str) -> str:
    """Generate a stable description for photo entity in KG.

    Only includes immutable attributes (file name, date from stem).
    Person names are NOT included — they belong to person entities and
    are expressed via edges (features, co_occurs_with).
    This prevents the description from becoming stale when a person is renamed.
    """
    # Extract date from stem if possible (format: YYYYMMDD_HHMMSS)
    date_part = ""
    if len(normalized_stem) >= 8 and normalized_stem[:8].isdigit():
        try:
            from datetime import datetime
            dt = datetime.strptime(normalized_stem[:8], "%Y%m%d")
            # Validate: month 1-12, day 1-31 (strptime already rejects invalid dates)
            if dt.year < 1900 or dt.year > 2100:
                raise ValueError("year out of reasonable range")
            date_part = dt.strftime("%Y年%m月%d日")
        except ValueError:
            pass

    parts = [f"照片 {normalized_stem}"]
    if date_part:
        parts.append(f"拍摄于{date_part}")

    return "，".join(parts)


def format_photo_ingest_data(
    file_path: str, abstract: str, detected_persons: list
) -> dict:
    """格式化照片信息为结构化实体+关系，供 inject_custom_kg 注入。

    返回 {"entities": [...], "relationships": [...]}，不触发 LLM。

    实体命名原则：完全遵循 LightRAG 的自然语言命名体系，不使用冒号前缀。
    - 照片实体名：使用文件名 stem（如 20090603_092316），确保稳定不随 abstract 变化
    - 人物实体名：使用真名或 auto_label，如"任飞"、"未命名人物_1"
    - 禁止 photo:{stem}、person:{uuid} 等编程风格命名

    稳定性保证：照片实体名 = 文件名 stem，description 不含人名，不受人物改名影响。
    人物信息完全通过边（features, co_occurs_with）表达。

    路径统一归一化（正斜杠 + 小写），确保 file_path 不会因大小写或斜杠方向差异而分裂。
    """
    # 归一化路径：统一为正斜杠 + 小写，避免大小写分裂
    normalized_path = file_path.replace("\\", "/").lower()
    normalized_stem = Path(normalized_path).stem  # 归一化后的文件名 stem

    # 照片实体名：使用文件名 stem，确保稳定不随 abstract 变化
    # abstract 仅放入 description，用于搜索和语义理解
    photo_entity_name = normalized_stem

    # 照片实体
    entities = [
        {
            "entity_name": photo_entity_name,
            "entity_type": "Photo",
            "description": _generate_stable_description(normalized_stem, abstract),
            "file_path": normalized_path,
            "source_id": normalized_path,
        }
    ]

    relationships = []

    # 注意：照片和人物实体不直接连接 Niu 根节点。
    # 照片通过 photo → person (features) 关系自然可达。
    # 脑区归属由 brain_region 系统自动管理。

    # 人物实体 + 关系
    person_names = []
    seen_person_entity_names: set[str] = set()  # deduplicate: same person may appear via multiple faces
    for p in detected_persons:
        pname = p.get("name", "")
        auto_label = p.get("auto_label", "")
        # 已命名用真名，未命名用 auto_label
        # 未命名判定：name 为空 / name 以"未命名人物"开头 / name 与 auto_label 相同
        is_unnamed = not pname or pname.startswith("未命名人物") or pname == auto_label
        entity_name = auto_label if is_unnamed else pname
        if not entity_name:
            continue

        # Deduplicate: skip if this person entity_name was already added
        if entity_name in seen_person_entity_names:
            continue
        seen_person_entity_names.add(entity_name)

        person_names.append(entity_name)
        # 保留 person_id（UUID）到 description 中，便于 KG 实体溯源
        person_pid = p.get("id", "")
        desc = f"{entity_name}，出现在照片{normalized_stem}中"
        if person_pid:
            desc += f"（person_id={person_pid}）"
        entities.append({
            "entity_name": entity_name,
            "entity_type": "person",
            "description": desc,
            "file_path": normalized_path,
            "source_id": normalized_path,  # 与 chunk 的 source_id 一致，确保映射成功
        })

        # 照片 → 人物 features 边
        relationships.append({
            "src_id": photo_entity_name,
            "tgt_id": entity_name,
            "keywords": "features",
            
            "file_path": normalized_path,
            "source_id": normalized_path,  # 与 chunk 的 source_id 一致，确保映射成功
        })

        # 注意：人物实体不直接连接 Niu 根节点。
        # 人物通过 photo → person (features) 关系自然可达。
        # 脑区归属由 brain_region 系统自动管理。

    # 多人同框：co_occurs_with 双向关系
    # Safety: deduplicate person_names to prevent spurious self-referencing edges
    person_names = list(dict.fromkeys(person_names))
    for i in range(len(person_names)):
        for j in range(i + 1, len(person_names)):
            a, b = person_names[i], person_names[j]
            relationships.append({
                "src_id": a,
                "tgt_id": b,
                "keywords": "co_occurs_with",
                
                "file_path": normalized_path,
                "source_id": normalized_path,  # 与 chunk 的 source_id 一致，确保映射成功
            })

    return {"entities": entities, "relationships": relationships}


def sync_photo_to_kg(file_path: str, abstract: str, detected_persons: list, force: bool = False) -> dict:
    """同步照片信息到知识图谱（结构化注入）

    流程：
    1. custom_kg(entities + relationships + chunks) -- 注入实体+关系（带chunk关联source_id）
    2. 标记 kg_synced

    注意：ainsert (Step 2) 已禁用，因为它会阻塞 LightRAG 事件循环长达 90+ 秒。
    仅保留 Step 1 的结构化注入，足以支持基本的图谱查询。

    防重复：如果照片已标记 kg_synced=1 且 force=False，跳过整个流程。
    name_person 改名后不应重新 sync_photo_to_kg，只更新人物实体本身。

    Args:
        file_path: 照片文件路径
        abstract: 照片摘要
        detected_persons: 检测到的人物列表
        force: 是否强制重新同步
    """
    # 防重复检查：已 kg_synced 的照片不重新注入（除非 force=True）
    if not force:
        try:
            conn = get_connection()
            cursor = conn.execute(
                "SELECT kg_synced FROM photos WHERE file_path = ?",
                (file_path,),
            )
            row = cursor.fetchone()
            if row and row[0] == 1:
                logger.info(f"[KG] Photo {file_path} already synced to KG, skipping")
                return {"status": "skipped", "reason": "already synced", "kg_entities": []}
        except Exception as e:
            logger.warning(f"[KG] kg_synced check failed for {file_path}: {e}")

    # 直接同步执行（不再使用异步模式，避免阻塞 LightRAG 事件循环）
    return _do_sync_photo_to_kg_sync(file_path, abstract, detected_persons)


def _do_sync_photo_to_kg_sync(file_path: str, abstract: str, detected_persons: list) -> dict:
    """同步执行 KG 同步（内部函数）"""
    try:
        from agent.tool_registry import get_registry

        data = format_photo_ingest_data(file_path, abstract, detected_persons)
        normalized_path = file_path.replace("\\", "/").lower()
        normalized_stem = Path(normalized_path).stem
        registry = get_registry()

        # --- 构建 chunk_text（Step 2 用，明确引用实体名让 LLM 能识别） ---
        # 照片实体名 = 文件名 stem（如"20090603_092316"），稳定不变；
        # 人物实体名 = 自然语言名（如"任飞"、"未命名人物_1"）
        # LLM 在 ainsert 时能识别这些名称并建立语义边
        entity_names = [e["entity_name"] for e in data["entities"]]
        chunk_text = (
            f"照片 {normalized_stem}：{abstract}\n"
            f"实体：{', '.join(entity_names)}\n"
        )
        person_list = ", ".join(
            e["entity_name"] for e in data["entities"]
            if e.get("entity_type") == "person"
        )
        if person_list:
            chunk_text += f"人物：{person_list}\n"

        custom_kg_fn = registry.get("lightrag-server/lightrag_insert_custom_kg")
        if not custom_kg_fn:
            logger.warning("[KG] lightrag_insert_custom_kg not available in registry")
            return {"status": "error", "reason": "lightrag_insert_custom_kg not available", "kg_entities": []}

        # --- Step 1: 注入实体 + 关系（照片 + 人物 + features），带 chunk 关联 source_id ---
        entity_result = custom_kg_fn(
            entities=data["entities"],
            relationships=data["relationships"],
            chunks=[{
                "content": chunk_text,
                "source_id": normalized_path,
                "file_path": normalized_path,
            }],
            source_id=normalized_path,
        )
        if not entity_result or entity_result.get("status") != "ok":
            logger.warning(f"[KG] Step1 entity+relationship injection failed for {file_path}: {entity_result}")
            return {"status": "error", "reason": f"Step1 entity+relationship injection failed: {entity_result}", "kg_entities": []}
        logger.info(f"[KG] Step1 ok: {len(data['entities'])} entities + {len(data['relationships'])} relationships injected for {normalized_stem}")

        # --- Step 2: ainsert 让 LLM 处理文本，建立语义连接 ---
        # DISABLED: ainsert 会阻塞 LightRAG 事件循环长达 90+ 秒，
        # 导致后续的 query 操作超时。暂时禁用，仅保留 Step 1 的结构化注入。
        # Step 1 已经注入了实体和关系，足以支持基本的图谱查询。
        # TODO: 考虑使用独立的队列或进程来处理 ainsert。
        #
        # chunk_text 中明确引用了 Step 1 的实体名（自然语言），LLM 能识别并合并
        # insert_fn = registry.get("lightrag-server/lightrag_insert")
        # if insert_fn:
        #     try:
        #         insert_result = insert_fn(
        #             content=chunk_text,
        #             file_path=normalized_path,  # 避免 unknown_source
        #             doc_id=f"doc-{normalized_stem}",
        #         )
        #         if insert_result and insert_result.get("status") == "ok":
        #             logger.info(f"[KG] Step2 ok: ainsert completed for {normalized_stem}")
        #         else:
        #             logger.warning(f"[KG] Step2 ainsert returned non-ok for {file_path}: {insert_result}")
        #     except Exception as e:
        #         logger.warning(f"[KG] Step2 ainsert failed for {file_path}: {e}")
        # else:
        #     logger.warning("[KG] lightrag_insert not available in registry, skipping Step2 ainsert")
        logger.info(f"[KG] Step2 ainsert skipped (disabled to prevent event loop blocking)")

        # --- Step 3: 标记 kg_synced ---
        # Only mark as KG-synced if Step 1 succeeded (Step2 failure is non-fatal)
        try:
            with _db_write_lock:
                conn = get_connection()
                conn.execute(
                    "UPDATE photos SET kg_synced = 1 WHERE file_path = ?",
                    (file_path,),
                )
                conn.commit()
        except Exception as e:
            logger.warning(f"[KG] Failed to mark kg_synced for {file_path}: {e}")

        # Build entity list for return value (so downstream agents know what's in the graph)
        kg_entities = [
            {"entity_name": e["entity_name"], "entity_type": e["entity_type"]}
            for e in data["entities"]
        ]

        return {"status": "success", "doc_uri": file_path, "kg_entities": kg_entities}

    except Exception as e:
        logger.warning(f"[KG] Photo sync failed: {e}")
        return {"status": "error", "reason": str(e), "kg_entities": []}



def sync_video_to_kg(file_path: str, abstract: str) -> dict:
    """同步视频到 LightRAG 知识图谱。

    视频描述通过 ainsert() 入库，实体提取由 LightRAG 自动完成。
    """
    try:
        from niu_api.internal.lightrag_manager import get_lightrag, call_async

        rag = get_lightrag()
        if rag is None:
            logger.warning("[KG] LightRAG not available, skipping video KG sync")
            return {"status": "skipped", "reason": "LightRAG not available"}

        content = f"[Video: {file_path}]\n{abstract}"
        call_async(rag.ainsert(content), timeout=600)
        logger.info(f"[KG] Video ingested into LightRAG: {file_path}")
        return {"status": "success", "doc_uri": file_path}

    except Exception as e:
        logger.warning(f"[KG] Video sync failed: {e}")
        return {"status": "error", "reason": str(e)}


# ============== 模型路径（人脸识别用） ==============


def _detect_available_providers() -> list[str]:
    """
    自动检测可用的 ONNX Runtime ExecutionProvider

    优先级：CUDA > DirectML > CPU

    Returns:
        可用的 provider 列表
    """
    try:
        import onnxruntime as ort

        available = ort.get_available_providers()
        logger.info(f"[GPU_DETECT] Available providers: {available}")

        # 按优先级排序
        priority = ["CUDAExecutionProvider", "DmlExecutionProvider", "CPUExecutionProvider"]

        selected = []
        for provider in priority:
            if provider in available:
                selected.append(provider)

        # 如果没有找到任何优先 provider，使用所有可用的
        if not selected:
            selected = available

        logger.info(f"[GPU_DETECT] Selected providers: {selected}")
        return selected

    except Exception as e:
        logger.warning(f"[GPU_DETECT] Failed to detect providers: {e}, using CPU")
        return ["CPUExecutionProvider"]


def get_models_dir() -> Path:
    """获取模型目录路径（人脸识别模型使用）"""
    if "NIU_MODELS_PATH" in os.environ:
        return Path(os.environ["NIU_MODELS_PATH"])
    return Path(__file__).parent.parent.parent.parent.parent / "models"


# ============== Face recognition model ==============
_face_model = None
_last_model_use_time = None
_model_in_use = False  # 防止卸载定时器在推理期间卸载模型
_model_lock = threading.Lock()  # 保护 _model_in_use 和模型卸载的临界区
MODEL_IDLE_TIMEOUT_SECONDS = 300  # 5 分钟无使用自动卸载
_model_check_interval = 60  # 每 60 秒检查一次


def _start_model_unload_timer():
    """启动后台定时器，定期检查并卸载空闲模型"""
    import time

    def check_and_unload():
        global _face_model, _last_model_use_time
        while True:
            time.sleep(_model_check_interval)

            with _model_lock:
                if _face_model is not None and _last_model_use_time is not None and not _model_in_use:
                    idle_seconds = (datetime.now() - _last_model_use_time).total_seconds()
                    if idle_seconds > MODEL_IDLE_TIMEOUT_SECONDS:
                        logger.info(
                            f"[MODEL_UNLOAD] Model idle for {idle_seconds:.0f}s, unloading ~326MB..."
                        )
                        _face_model = None
                        _last_model_use_time = None
                        # 不调用 gc.collect()，让 Python 自然回收
                        # 避免 detect_faces() 正在使用模型时被释放
                        logger.info("[MODEL_UNLOAD] Face model unloaded, memory released")

    thread = threading.Thread(target=check_and_unload, daemon=True)
    thread.start()
    logger.info("[MODEL_UNLOAD] Background unload timer started")


def get_face_model():
    """Get or load InsightFace model. Local first, download if not found.

    自动检测 GPU，优先使用 GPU 加速，无 GPU 时降级到 CPU。
    """
    import sys

    global _face_model, _last_model_use_time

    if _face_model is None:
        try:
            print(
                "[GET_FACE_MODEL] Starting to load InsightFace...",
                file=sys.stderr,
                flush=True,
            )
            logger.info("[GET_FACE_MODEL] Starting to load InsightFace...")
            from insightface.app import FaceAnalysis

            # 本地路径
            models_dir = get_models_dir()
            print(
                f"[GET_FACE_MODEL] Models dir: {models_dir}",
                file=sys.stderr,
                flush=True,
            )
            logger.info(f"[GET_FACE_MODEL] Models dir: {models_dir}")

            # 检查本地是否存在
            local_model_path = models_dir / "models" / "buffalo_l"

            if local_model_path.exists():
                print(
                    f"[GET_FACE_MODEL] Local model exists: {local_model_path}",
                    file=sys.stderr,
                    flush=True,
                )
                logger.info(
                    f"[GET_FACE_MODEL] Loading face model from local: {local_model_path}"
                )
            else:
                # 本地没有，需要下载
                print(
                    f"[GET_FACE_MODEL] Local model not found, will download...",
                    file=sys.stderr,
                    flush=True,
                )
                logger.info(
                    f"[GET_FACE_MODEL] Local face model not found, will download..."
                )
                logger.info(f"[GET_FACE_MODEL] Expected path: {local_model_path}")

            print(
                "[GET_FACE_MODEL] Creating FaceAnalysis instance...",
                file=sys.stderr,
                flush=True,
            )
            logger.info("[GET_FACE_MODEL] Creating FaceAnalysis instance...")

            # 自动检测可用的 ExecutionProvider
            providers = _detect_available_providers()

            # 关键：临时抑制 stdout，防止 ONNX Runtime 污染 MCP stdio 通信
            import os
            import contextlib

            @contextlib.contextmanager
            def suppress_stdout():
                """临时抑制 stdout（ONNX Runtime 会污染 stdout）"""
                devnull = open(os.devnull, 'w')
                old_stdout = os.dup(1)
                try:
                    os.dup2(devnull.fileno(), 1)
                    yield
                finally:
                    os.dup2(old_stdout, 1)
                    devnull.close()
                    os.close(old_stdout)

            with suppress_stdout():
                _face_model = FaceAnalysis(
                    name="buffalo_l",
                    root=str(models_dir),
                    providers=providers,
                )
                # ctx_id: 0 = GPU, -1 = CPU
                ctx_id = 0 if "CUDAExecutionProvider" in providers else -1
                _face_model.prepare(ctx_id=ctx_id)

            print(
                f"[GET_FACE_MODEL] FaceAnalysis created with providers: {providers}",
                file=sys.stderr,
                flush=True,
            )
            logger.info(f"[GET_FACE_MODEL] FaceAnalysis created with providers: {providers}")

            print(
                f"[GET_FACE_MODEL] Face model loaded: buffalo_l (ctx_id={ctx_id})",
                file=sys.stderr,
                flush=True,
            )
            logger.info(f"[GET_FACE_MODEL] Face model loaded: buffalo_l (ctx_id={ctx_id})")
        except ImportError as e:
            print(
                f"[GET_FACE_MODEL] insightface not installed: {e}",
                file=sys.stderr,
                flush=True,
            )
            logger.warning(f"[GET_FACE_MODEL] insightface not installed: {e}")
            return None
        except Exception as e:
            print(
                f"[GET_FACE_MODEL] Failed to load face model: {e}",
                file=sys.stderr,
                flush=True,
            )
            logger.warning(f"[GET_FACE_MODEL] Failed to load face model: {e}")
            return None
    else:
        print("[GET_FACE_MODEL] Using cached face model", file=sys.stderr, flush=True)
        logger.info("[GET_FACE_MODEL] Using cached face model")

    # 更新最后使用时间
    _last_model_use_time = datetime.now()
    return _face_model


def unload_face_model():
    """卸载人脸识别模型，释放内存"""
    global _face_model, _last_model_use_time
    if _face_model is not None:
        _face_model = None
        _last_model_use_time = None
        # 不调用 gc.collect()，让 Python 自然回收
        # 避免 detect_faces() 正在使用模型时被释放
        logger.info("Face model unloaded, ~326MB memory released")
    return {"status": "success"}


def update_co_occurrences(persons: list[dict], photo_taken_at: str | None = None):
    """更新同框关系"""
    if len(persons) < 2:
        return

    try:
        conn = get_connection()
        now = photo_taken_at or datetime.now().isoformat()

        # 每对人创建/更新同框关系
        for i, person_a in enumerate(persons):
            for person_b in persons[i + 1 :]:
                # 确保 a < b（避免重复）
                a_id, b_id = sorted([person_a["id"], person_b["id"]])

                # 检查是否已有关系
                cursor = conn.execute(
                    "SELECT count FROM co_occurrences WHERE person_a_id = ? AND person_b_id = ?",
                    (a_id, b_id),
                )
                row = cursor.fetchone()

                if row:
                    # 更新
                    conn.execute(
                        """UPDATE co_occurrences 
                           SET count = count + 1, last_seen = ? 
                           WHERE person_a_id = ? AND person_b_id = ?""",
                        (now, a_id, b_id),
                    )
                else:
                    # 创建
                    conn.execute(
                        """INSERT INTO co_occurrences (person_a_id, person_b_id, count, first_seen, last_seen)
                           VALUES (?, ?, 1, ?, ?)""",
                        (a_id, b_id, now, now),
                    )

        # 不在此处 commit，由调用者控制事务原子性
        logger.info(f"[CO_OCCURRENCE] Updated {len(persons)} persons co-occurrence")
    except Exception as e:
        logger.warning(f"[CO_OCCURRENCE] Failed: {e}")


# ============== 配置路径 ==============


def get_config_path() -> Path:
    """获取配置目录路径"""
    if "NIU_CONFIG_PATH" in os.environ:
        return Path(os.environ["NIU_CONFIG_PATH"])
    return Path.home() / ".niu"


def get_memory() -> dict:
    """读取 memory.json"""
    memory_path = get_config_path() / "memory.json"
    if memory_path.exists():
        with open(memory_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {"workspace": {"path": str(Path.home() / "Documents" / "niu")}}


def get_workspace_path() -> Path:
    """获取工作区路径（3 层优先级，与 resolve_vector_db_path 一致）

    Priority:
    1. WORKSPACE_PATH 环境变量（由 Go 启动器 main.go 设置）
    2. ~/.niu/memory.json 的 workspace.path

    Auto-creates workspace directory if it doesn't exist.
    """
    # 1. WORKSPACE_PATH 环境变量
    if "WORKSPACE_PATH" in os.environ:
        ws = Path(os.environ["WORKSPACE_PATH"])
        if not ws.is_absolute():
            raise ValueError(f"WORKSPACE_PATH 必须是绝对路径: {ws}")
        # 解析符号链接
        ws = ws.resolve()
        if ws.exists() and not ws.is_dir():
            raise ValueError(f"WORKSPACE_PATH 指向的不是目录: {ws}")
        try:
            ws.mkdir(parents=True, exist_ok=True)
        except PermissionError as e:
            raise ValueError(f"无权限创建工作区目录: {ws}。请检查权限。") from e
        except FileExistsError as e:
            raise ValueError(f"路径已存在但不是目录（可能是文件）: {ws}") from e
        return ws
    # 2. 从 ~/.niu/memory.json 读取 workspace.path
    memory_path = Path.home() / ".niu" / "memory.json"
    if memory_path.exists():
        try:
            with open(memory_path, "r", encoding="utf-8") as f:
                memory = json.load(f)
            workspace_path = memory.get("workspace", {}).get("path")
            if workspace_path and workspace_path.strip():
                ws = Path(workspace_path)
                if not ws.is_absolute():
                    raise ValueError(f"workspace.path 必须是绝对路径: {workspace_path}")
                # 解析符号链接
                ws = ws.resolve()
                if ws.exists() and not ws.is_dir():
                    raise ValueError(f"workspace.path 指向的不是目录: {workspace_path}")
                try:
                    ws.mkdir(parents=True, exist_ok=True)
                except PermissionError as e:
                    raise ValueError(f"无权限创建工作区目录: {ws}。请检查权限。") from e
                except FileExistsError as e:
                    raise ValueError(f"路径已存在但不是目录（可能是文件）: {ws}") from e
                return ws
        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"无法从 {memory_path} 解析 JSON: {e}。请检查 memory.json 格式是否正确。") from e
    raise ValueError(
        f"无法确定工作区路径：~/.niu/memory.json 不存在或缺少 workspace.path 配置。"
        f"请在 ~/.niu/memory.json 中设置 workspace.path，或设置 WORKSPACE_PATH 环境变量。"
    )


def get_preferences() -> dict:
    """读取 preferences.json"""
    prefs_path = get_config_path() / "preferences.json"
    if prefs_path.exists():
        with open(prefs_path, "r", encoding="utf-8") as f:
            return json.load(f)
    raise FileNotFoundError(f"preferences.json not found at {prefs_path}")


# ============== EXIF Extraction ==============


def extract_exif(file_path: str) -> dict:
    """Extract EXIF data from photo."""
    result: dict[str, Any] = {
        "taken_at": None,
        "location": None,
        "camera": None,
    }

    try:
        from PIL import Image
        from PIL.ExifTags import TAGS, GPSTAGS

        img = Image.open(file_path)
        exif_data = img._getexif()  # type: ignore

        if not exif_data:
            return result

        # Extract datetime
        for tag_id, value in exif_data.items():
            tag = TAGS.get(tag_id, tag_id)

            if tag == "DateTimeOriginal":
                result["taken_at"] = value
            elif tag == "DateTimeDigitized":
                if not result["taken_at"]:
                    result["taken_at"] = value
            elif tag == "Model":
                result["camera"] = value
            elif tag == "Make":
                if result["camera"]:
                    result["camera"] = f"{value} {result['camera']}"
                else:
                    result["camera"] = value

        # Extract GPS
        for tag_id, value in exif_data.items():
            tag = TAGS.get(tag_id, tag_id)
            if tag == "GPSInfo":
                gps_data: dict[str, Any] = {}
                for gps_tag, gps_val in value.items():
                    gps_tag_name = str(GPSTAGS.get(gps_tag, gps_tag))
                    gps_data[gps_tag_name] = gps_val

                lat = gps_data.get("GPSLatitude")
                lat_ref = gps_data.get("GPSLatitudeRef")
                lon = gps_data.get("GPSLongitude")
                lon_ref = gps_data.get("GPSLongitudeRef")

                if lat and lat_ref and lon and lon_ref:
                    lat_val = lat[0] + lat[1] / 60 + lat[2] / 3600
                    if lat_ref == "S":
                        lat_val = -lat_val
                    lon_val = lon[0] + lon[1] / 60 + lon[2] / 3600
                    if lon_ref == "W":
                        lon_val = -lon_val
                    result["location"] = f"{lat_val:.6f},{lon_val:.6f}"

                break

    except ImportError:
        logger.warning("PIL not installed, EXIF extraction disabled")
    except Exception as e:
        logger.warning(f"EXIF extraction failed: {e}")

    return result


# ============== Face Recognition ==============


def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    """Calculate cosine similarity between two vectors."""
    norm_a = np.linalg.norm(a)
    norm_b = np.linalg.norm(b)
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return float(np.dot(a, b) / (norm_a * norm_b))


def get_all_persons() -> list[dict]:
    """Get all persons from database."""
    conn = get_connection()
    cursor = conn.execute(
        "SELECT id, name, auto_label, center_embedding, threshold_adjustment, photo_count FROM persons"
    )
    persons = []
    for row in cursor.fetchall():
        persons.append(
            {
                "id": row[0],
                "name": row[1],
                "auto_label": row[2],
                "center_embedding": np.frombuffer(row[3], dtype=np.float32)
                if row[3]
                else None,
                "threshold_adjustment": row[4] or 0,
                "photo_count": row[5] or 0,
            }
        )
    return persons


def get_next_auto_label() -> str:
    """Get next auto label for unnamed person."""
    conn = get_connection()
    cursor = conn.execute(
        "SELECT auto_label FROM persons WHERE auto_label LIKE '未命名人物_%' ORDER BY auto_label DESC LIMIT 1"
    )
    row = cursor.fetchone()
    if row:
        # Extract number and increment
        try:
            num = int(row[0].split("_")[-1])
            return f"未命名人物_{num + 1}"
        except (ValueError, IndexError, AttributeError):
            pass
    return "未命名人物_1"


def search_persons(query: str, limit: int = 10) -> list[dict]:
    """
    搜索人物（按名字模糊匹配）

    Args:
        query: 搜索词（人名）
        limit: 返回数量

    Returns:
        匹配的人物列表
    """
    try:
        conn = get_connection()
        # 使用 SQL LIKE 模糊匹配，转义 LIKE 通配符
        escaped = query.replace("%", "\\%").replace("_", "\\_")
        cursor = conn.execute(
            """SELECT id, name, auto_label, photo_count
               FROM persons
               WHERE name IS NOT NULL AND name LIKE ? ESCAPE '\\'""",
            (f"%{escaped}%",),
        )

        results = []
        for row in cursor.fetchall():
            person_id, name, auto_label, photo_count = row
            results.append(
                {
                    "id": person_id,
                    "name": name,
                    "auto_label": auto_label,
                    "photo_count": photo_count or 0,
                }
            )

        return results[:limit]

    except Exception as e:
        logger.exception(f"[SEARCH_PERSONS] Failed: {e}")
        return []


def get_unnamed_persons() -> list[dict]:
    """获取未命名人物列表，包含多张代表照片和人脸框

    返回每个人的多张照片（最多3张），供主 Agent 轮流展示。
    只返回存在的照片文件。如果人物没有任何有效照片，标记 has_valid_photos=false。
    """
    from pathlib import Path

    conn = get_connection()
    cursor = conn.execute(
        """SELECT id, name, auto_label, photo_count, first_seen, last_seen 
           FROM persons 
           WHERE name IS NULL 
           ORDER BY photo_count DESC"""
    )

    persons = []
    for row in cursor.fetchall():
        person_id = row[0]

        # 获取该人物的所有照片（去重，同一照片只取第一个bbox）
        photo_cursor = conn.execute(
            """SELECT p.file_path, f.bounding_box, p.id
               FROM photos p
               JOIN faces f ON f.photo_id = p.id
               WHERE f.person_id = ?
               ORDER BY p.taken_at DESC""",
            (person_id,),
        )

        seen_photo_ids = set()
        photos = []
        for photo_row in photo_cursor.fetchall():
            # 同一照片只取一次（避免多个人脸导致重复）
            photo_id = photo_row[2]
            if photo_id in seen_photo_ids:
                continue
            seen_photo_ids.add(photo_id)

            # 检查照片文件是否存在
            if not Path(photo_row[0]).exists():
                continue

            bbox = None
            if photo_row[1]:
                try:
                    bbox = json.loads(photo_row[1])
                except (json.JSONDecodeError, TypeError, ValueError):
                    logger.warning(f"[UnnamedPersons] Invalid bbox JSON for photo {photo_row[0]}")

            # 在原图上画人脸红框
            boxed_path = None
            if bbox:
                boxed_path = draw_face_boxes_on_image(photo_row[0], [bbox])

            photos.append(
                {
                    "file_path": photo_row[0],
                    "boxed_path": boxed_path,  # 带红框图路径
                }
            )

            # 最多返回3张存在的照片
            if len(photos) >= 3:
                break

        persons.append(
            {
                "id": person_id,
                "name": row[1],
                "auto_label": row[2],
                "photo_count": row[3] or 0,
                "first_seen": row[4],
                "last_seen": row[5],
                "photos": photos,  # 只包含存在的照片
                "has_valid_photos": len(photos) > 0,  # 是否有有效照片
            }
        )

    return persons


def delete_person(person_id: str) -> dict:
    """删除人物及其所有关联数据

    警告：这会删除人物图谱中的节点，请谨慎使用。
    只有在用户明确要求删除时才调用。
    """
    conn = get_connection()

    # 检查人物是否存在
    cursor = conn.execute(
        "SELECT name, auto_label FROM persons WHERE id = ?", (person_id,)
    )
    row = cursor.fetchone()
    if not row:
        return {"status": "error", "message": f"Person not found: {person_id}"}

    person_name = row[0] if row[0] else row[1]

    # 找出同照片中其他人物（删除 faces 前查询）
    co_person_ids = conn.execute(
        """SELECT DISTINCT f2.person_id FROM faces f1
           JOIN faces f2 ON f1.photo_id = f2.photo_id
           WHERE f1.person_id = ? AND f2.person_id != ? AND f2.person_id IS NOT NULL""",
        (person_id, person_id),
    ).fetchall()

    # 写事务：删除关联数据并提交
    with _db_write_lock:
        # 删除关联的人脸记录
        conn.execute("DELETE FROM faces WHERE person_id = ?", (person_id,))

        # 重新计算同照片中其他人物的 photo_count
        for (pid,) in co_person_ids:
            actual_count = conn.execute(
                "SELECT COUNT(DISTINCT photo_id) FROM faces WHERE person_id = ?",
                (pid,),
            ).fetchone()[0]
            conn.execute(
                "UPDATE persons SET photo_count = ? WHERE id = ?",
                (actual_count, pid),
            )

        # 删除同框关系
        conn.execute(
            "DELETE FROM co_occurrences WHERE person_a_id = ? OR person_b_id = ?",
            (person_id, person_id),
        )

        # 删除人物记录
        conn.execute("DELETE FROM persons WHERE id = ?", (person_id,))

        conn.commit()

    # 同步删除知识图谱中的实体
    if person_name:
        try:
            from agent.tool_registry import get_registry
            registry = get_registry()
            delete_fn = registry.get("lightrag-server/lightrag_delete_entity")
            if delete_fn:
                delete_fn(entity_name=person_name)
                logger.info(f"[DELETE_PERSON] KG entity deleted: {person_name}")
            else:
                logger.warning("[DELETE_PERSON] lightrag_delete_entity not available in registry")
        except Exception as e:
            logger.warning(f"[DELETE_PERSON] LightRAG entity deletion failed: {e}")
    else:
        logger.warning(f"[DELETE_PERSON] No entity name for person {person_id}, skipping KG deletion")

    return {
        "status": "success",
        "message": f"Deleted person: {person_name}",
        "person_id": person_id,
    }


def cleanup_deleted_photos() -> dict:
    """清理已删除照片的数据库记录

    扫描 photos 表，删除文件不存在的照片记录及其关联的 faces 记录。
    在删除照片文件后调用此工具清理数据库。

    返回:
    - deleted_photos: 删除的照片记录数
    - deleted_faces: 删除的人脸记录数
    """
    from pathlib import Path

    conn = get_connection()

    # 找出文件不存在的照片
    cursor = conn.execute("SELECT id, file_path FROM photos")
    deleted_photo_ids = []

    for row in cursor.fetchall():
        photo_id = row[0]
        file_path = row[1]
        if not Path(file_path).exists():
            deleted_photo_ids.append(photo_id)

    if not deleted_photo_ids:
        return {
            "status": "success",
            "message": "No deleted photos to cleanup",
            "deleted_photos": 0,
            "deleted_faces": 0,
        }

    placeholders = ",".join("?" * len(deleted_photo_ids))

    # 找出受影响的 person_id（必须在删除 faces 之前查询）
    affected_person_ids = conn.execute(
        f"SELECT DISTINCT person_id FROM faces WHERE photo_id IN ({placeholders}) AND person_id IS NOT NULL",
        deleted_photo_ids,
    ).fetchall()

    # 写事务：删除关联记录并提交
    with _db_write_lock:
        # 删除关联的 faces 记录
        cursor = conn.execute(
            f"DELETE FROM faces WHERE photo_id IN ({placeholders})",
            deleted_photo_ids,
        )
        deleted_faces = cursor.rowcount

        # 删除 photos 记录
        cursor = conn.execute(
            f"DELETE FROM photos WHERE id IN ({placeholders})",
            deleted_photo_ids,
        )
        deleted_photos = cursor.rowcount

        # 重新计算受影响人物的 photo_count
        for (pid,) in affected_person_ids:
            actual_count = conn.execute(
                "SELECT COUNT(DISTINCT photo_id) FROM faces WHERE person_id = ?",
                (pid,),
            ).fetchone()[0]
            conn.execute(
                "UPDATE persons SET photo_count = ? WHERE id = ?",
                (actual_count, pid),
            )

        conn.commit()

    return {
        "status": "success",
        "message": f"Cleaned up {deleted_photos} deleted photos, {deleted_faces} faces",
        "deleted_photos": deleted_photos,
        "deleted_faces": deleted_faces,
    }


def draw_face_boxes_on_image(file_path: str, bbox_list: list[list[float]]) -> str | None:
    """在原图上画人脸红框，保存到临时目录，返回带框图路径。

    Args:
        file_path: 原图路径
        bbox_list: 人脸框列表，每个为 [x1, y1, x2, y2]

    Returns:
        带框图路径（~/.niu/tmp/ 下），失败返回 None
    """
    try:
        import cv2

        # 读取图片（支持中文路径）
        with open(file_path, "rb") as f:
            img_bytes = np.frombuffer(f.read(), dtype=np.uint8)
        img = cv2.imdecode(img_bytes, cv2.IMREAD_COLOR)
        if img is None:
            logger.warning(f"[DrawBox] Cannot read image: {file_path}")
            return None

        # 画红框
        for bbox in bbox_list:
            if len(bbox) == 4:
                x1, y1, x2, y2 = int(bbox[0]), int(bbox[1]), int(bbox[2]), int(bbox[3])
                cv2.rectangle(img, (x1, y1), (x2, y2), (0, 0, 255), 3)

        # 保存到临时目录（确定性文件名，相同输入复用同一文件）

        # 编码为 PNG bytes
        success, encoded = cv2.imencode(".png", img)
        if not success:
            logger.warning(f"[DrawBox] Failed to encode image: {file_path}")
            return None

        # 用路径+bbox+文件修改时间生成确定性文件名，避免重复生成
        bbox_key = "_".join(str(int(b)) for bbox in bbox_list for b in bbox)
        mtime = int(os.path.getmtime(file_path))
        name_hash = hashlib.md5(f"{file_path}:{bbox_key}:{mtime}".encode()).hexdigest()[:12]
        tmp_name = f"facebox_{name_hash}.png"

        # 获取临时目录（优先用 agent.tmp_dir，fallback 到 ~/.niu/tmp/）
        try:
            from agent.tmp_dir import get_tmp_dir
            tmp_dir = get_tmp_dir()
        except ImportError:
            tmp_dir = os.path.join(os.path.expanduser("~"), ".niu", "tmp")
            os.makedirs(tmp_dir, exist_ok=True)

        tmp_path = os.path.join(tmp_dir, tmp_name)
        if not os.path.exists(tmp_path):
            # 先写临时文件，再原子重命名，避免并发写或写失败导致损坏
            tmp_write_path = tmp_path + ".tmp"
            with open(tmp_write_path, "wb") as f:
                f.write(encoded.tobytes())
            os.replace(tmp_write_path, tmp_path)
        logger.info(f"[DrawBox] Saved boxed image to: {tmp_path}")
        return tmp_path

    except Exception as e:
        logger.error(f"[DrawBox] Error drawing face boxes: {e}")
        return None


def get_person_photos(person_id: str, limit: int = 5) -> dict:
    """获取某人物的所有照片（带人脸红框的图片路径）

    用于"换一张照片看看"的场景。
    在原图上画人脸红框，保存到临时目录，返回带框图路径。
    """
    conn = get_connection()

    # 验证人物存在
    cursor = conn.execute(
        "SELECT name, auto_label FROM persons WHERE id = ?", (person_id,)
    )
    row = cursor.fetchone()
    if not row:
        return {"status": "error", "message": f"Person not found: {person_id}"}

    person_name = row[0] if row[0] else row[1]

    # 获取该人物的所有照片（去重，同一照片只取第一个bbox）
    cursor = conn.execute(
        """SELECT p.file_path, f.bounding_box, p.taken_at, p.id
           FROM photos p
           JOIN faces f ON f.photo_id = p.id
           WHERE f.person_id = ?
           ORDER BY p.taken_at DESC""",
        (person_id,),
    )

    seen_photo_ids = set()
    photos = []
    for photo_row in cursor.fetchall():
        # 同一照片只取一次
        photo_id = photo_row[3]
        if photo_id in seen_photo_ids:
            continue
        seen_photo_ids.add(photo_id)

        if len(photos) >= limit:
            break

        bbox = None
        if photo_row[1]:
            try:
                bbox = json.loads(photo_row[1])
            except (json.JSONDecodeError, TypeError, ValueError):
                logger.warning(f"[PersonPhotos] Invalid bbox JSON for photo {photo_row[0]}")

        # 在原图上画人脸红框，保存到临时目录
        boxed_path = None
        if bbox and os.path.exists(photo_row[0]):
            boxed_path = draw_face_boxes_on_image(photo_row[0], [bbox])

        photos.append(
            {
                "file_path": photo_row[0],
                "boxed_path": boxed_path,  # 带红框图路径（~/.niu/tmp/），前端直接显示
                "taken_at": photo_row[2],
            }
        )

    return {
        "person_id": person_id,
        "person_name": person_name,
        "photos": photos,
    }


def match_face_to_person(
    face_embedding: np.ndarray, threshold: float = 0.7
) -> tuple[str | None, float]:
    """Match face embedding to existing person. Returns (person_id, similarity)."""
    persons = get_all_persons()

    best_match = None
    best_similarity = 0

    for person in persons:
        if person["center_embedding"] is None:
            continue

        similarity = cosine_similarity(face_embedding, person["center_embedding"])
        # Apply threshold adjustment (learning mechanism)
        adjusted_threshold = threshold - person["threshold_adjustment"]

        if similarity > adjusted_threshold and similarity > best_similarity:
            best_match = person["id"]
            best_similarity = similarity

    return best_match, best_similarity


def update_person_center(person_id: str, new_embedding: np.ndarray, increment_count: bool = True) -> None:
    """Update person center embedding with new face.

    Args:
        increment_count: If True, increment photo_count. Set to False when the same
                         person has multiple faces in the same photo (only count once).
    """
    conn = get_connection()

    # Get existing embedding and face count
    cursor = conn.execute(
        "SELECT center_embedding, photo_count FROM persons WHERE id = ?", (person_id,)
    )
    row = cursor.fetchone()

    if row and row[0]:
        existing = np.frombuffer(row[0], dtype=np.float32)
        photo_count = row[1] or 0

        # Incremental update: weighted average
        # existing center has weight photo_count, new embedding has weight 1
        updated = (existing * photo_count + new_embedding) / (photo_count + 1)
    else:
        updated = new_embedding

    # Update database
    if increment_count:
        conn.execute(
            "UPDATE persons SET center_embedding = ?, photo_count = photo_count + 1, last_seen = ? WHERE id = ?",
            (updated.tobytes(), datetime.now().isoformat(), person_id),
        )
    else:
        conn.execute(
            "UPDATE persons SET center_embedding = ?, last_seen = ? WHERE id = ?",
            (updated.tobytes(), datetime.now().isoformat(), person_id),
        )
    # 不在此处 commit，由调用者控制事务原子性


def detect_faces(file_path: str) -> list[dict]:
    """Detect faces in photo using InsightFace."""
    import sys

    print("[DETECT_FACES] Starting face detection...", file=sys.stderr, flush=True)

    logger.info(f"[DETECT_FACES] Starting face detection for: {file_path}")

    global _model_in_use
    with _model_lock:
        _model_in_use = True  # 在获取模型之前设置，防止卸载定时器在获取和使用之间卸载
    try:
        logger.info("[DETECT_FACES] Getting face model...")
        print("[DETECT_FACES] Getting face model...", file=sys.stderr, flush=True)
        face_model = get_face_model()
        if face_model is None:
            logger.warning("[DETECT_FACES] Face model not available")
            return []
        logger.info("[DETECT_FACES] Importing cv2...")
        print("[DETECT_FACES] Importing cv2...", file=sys.stderr, flush=True)
        import cv2
        import numpy as np

        # 使用 numpy 读取文件，解决中文路径问题
        # OpenCV 的 cv2.imread 不支持中文路径
        logger.info(f"[DETECT_FACES] Reading image file: {file_path}")
        print(f"[DETECT_FACES] Reading image file...", file=sys.stderr, flush=True)
        with open(file_path, "rb") as f:
            img_bytes = np.frombuffer(f.read(), dtype=np.uint8)

        logger.info("[DETECT_FACES] Decoding image with cv2...")
        print("[DETECT_FACES] Decoding image with cv2...", file=sys.stderr, flush=True)
        img = cv2.imdecode(img_bytes, cv2.IMREAD_COLOR)

        if img is None:
            logger.error(f"[DETECT_FACES] Cannot read image: {file_path}")
            return []

        logger.info(f"[DETECT_FACES] Image decoded, shape: {img.shape}")

        # Detect faces
        logger.info("[DETECT_FACES] Running face_model.get()...")
        print("[DETECT_FACES] Running face_model.get()...", file=sys.stderr, flush=True)
        faces = face_model.get(img)
        logger.info(f"[DETECT_FACES] Face detection complete, found {len(faces)} faces")

        results = []
        for face in faces:
            results.append(
                {
                    "embedding": face.embedding,
                    "bbox": face.bbox.tolist() if hasattr(face, "bbox") else [],
                    "confidence": float(face.det_score)
                    if hasattr(face, "det_score")
                    else 1.0,
                }
            )

        logger.info(f"[DETECT_FACES] Detected {len(results)} faces in {file_path}")
        return results

    except Exception as e:
        logger.exception(f"[DETECT_FACES] Face detection failed: {e}")
        return []
    finally:
        with _model_lock:
            _model_in_use = False


def generate_l0_abstract(person_names: list[str], taken_at: str | None) -> str:
    """Generate L0 abstract for photo."""
    parts = []

    if person_names:
        persons_str = "、".join(person_names)
        parts.append(f"{persons_str}合影")
    else:
        parts.append("单人照片")

    if taken_at:
        # Try to format the date nicely
        try:
            dt = datetime.strptime(taken_at, "%Y:%m:%d %H:%M:%S")
            parts.append(dt.strftime("%Y年%m月%d日"))
        except (ValueError, TypeError):
            parts.append(taken_at[:10] if len(taken_at) >= 10 else taken_at)

    return "，".join(parts)


# ============== Photo Tools ==============


def build_photo_storage_path(category: str, file_name: str) -> str:
    """构建照片存储路径（相对路径）"""
    prefs = get_preferences()
    structure = prefs["storage"]["structure"].get("photos", "{year}")

    now = datetime.now()
    path = structure.replace("{year}", str(now.year))
    path = path.replace("{month}", f"{now.month:02d}")
    path = path.replace("{date}", now.strftime("%Y-%m-%d"))
    path = path.replace("{category}", category)

    return path


def build_photo_file_name(source_name: str, taken_at: str | None = None) -> str:
    """构建照片文件名：日期_时间.扩展名

    不包含人物名，因为人物名会被修改，文件名应该只包含固定信息。
    - 有拍摄时间：20260402_092311.jpg
    - 无拍摄时间：20260402_143052.jpg（使用当前时间）
    """
    # 日期和时间
    if taken_at:
        try:
            # EXIF 格式: "2026:04:19 14:30:00" → 先尝试 strptime
            dt = datetime.strptime(taken_at, "%Y:%m:%d %H:%M:%S")
            date_str = dt.strftime("%Y%m%d")
            time_str = dt.strftime("%H%M%S")
        except (ValueError, TypeError):
            try:
                # ISO 格式: "2026-04-19T14:30:00" 或 "2026-04-19 14:30:00"
                dt = datetime.fromisoformat(taken_at.replace("Z", "+00:00"))
                date_str = dt.strftime("%Y%m%d")
                time_str = dt.strftime("%H%M%S")
            except (ValueError, TypeError, AttributeError):
                date_str = datetime.now().strftime("%Y%m%d")
                time_str = datetime.now().strftime("%H%M%S")
    else:
        date_str = datetime.now().strftime("%Y%m%d")
        time_str = datetime.now().strftime("%H%M%S")

    # 原始扩展名
    suffix = Path(source_name).suffix

    return f"{date_str}_{time_str}{suffix}"


def handle_photo_conflict(target_path: Path) -> str:
    """
    处理照片文件重名（防重名，不做相似度检测）
    返回最终存储路径
    """
    if not target_path.exists():
        return str(target_path)

    # 重名时自动改名
    stem = target_path.stem
    suffix = target_path.suffix
    parent = target_path.parent

    index = 1
    new_path = parent / f"{stem}_{index}{suffix}"
    while new_path.exists():
        index += 1
        new_path = parent / f"{stem}_{index}{suffix}"

    logger.info(f"[PHOTO_CONFLICT] 重名，改为: {new_path}")
    return str(new_path)


def ingest_photo(file_path: str, category: str | None = None, mode: str = "copy") -> dict:
    """Ingest photo with face detection and person matching."""
    conn = None
    final_path = None  # Track for cleanup on failure
    try:
        logger.info(f"[INGEST_PHOTO] Processing: {file_path}")

        source = Path(file_path)
        if not source.exists():
            return {
                "status": "error",
                "error_code": "FILE_NOT_FOUND",
                "message": f"File not found: {file_path}",
                "kg_entities": [],
            }

        # Get default category from preferences
        if category is None:
            try:
                prefs = get_preferences()
                category = prefs["categories"]["photos"][0]
            except (KeyError, TypeError, IndexError):
                category = "生活"

        # 确保 category 不为 None
        if category is None:
            category = "生活"

        # 1. Extract EXIF
        logger.info("[INGEST_PHOTO] Extracting EXIF...")
        exif = extract_exif(file_path)
        logger.info(f"[INGEST_PHOTO] EXIF: {exif}")

        # 2. Detect faces
        logger.info("[INGEST_PHOTO] Detecting faces...")
        faces = detect_faces(file_path)
        logger.info(f"[INGEST_PHOTO] Found {len(faces)} faces")

        # 3. Match/create persons
        detected_persons = []
        conn = get_connection()
        now = datetime.now().isoformat()

        try:
            with _db_write_lock:
                seen_person_ids = set()  # 防止同一人物多张人脸时 photo_count 重复递增
                for face_data in faces:
                    face_embedding = face_data["embedding"]

                    # Match to existing person
                    matched_id, similarity = match_face_to_person(face_embedding)

                    if matched_id:
                        person_id = matched_id
                        logger.info(
                            f"[INGEST_PHOTO] Matched face to person {person_id} (similarity: {similarity:.3f})"
                        )
                    else:
                        # Create new person
                        person_id = str(uuid.uuid4())
                        auto_label = get_next_auto_label()
                        similarity = 1.0  # 新人相似度设为 1.0（完全匹配）

                        conn.execute(
                            """INSERT INTO persons (id, auto_label, center_embedding, photo_count, first_seen, last_seen, created_at)
                               VALUES (?, ?, ?, 0, ?, ?, ?)""",
                            (person_id, auto_label, face_embedding.tobytes(), now, now, now),
                        )
                        logger.info(f"[INGEST_PHOTO] Created new person: {auto_label}")

                    # Update center embedding（仅首次出现时递增 photo_count）
                    is_new_photo = person_id not in seen_person_ids
                    update_person_center(person_id, face_embedding, increment_count=is_new_photo)
                    seen_person_ids.add(person_id)

                    # Get person info for response
                    cursor = conn.execute(
                        "SELECT name, auto_label FROM persons WHERE id = ?", (person_id,)
                    )
                    row = cursor.fetchone()
                    if row:
                        person_name = row[0] if row[0] else row[1]
                        detected_persons.append(
                            {
                                "id": person_id,
                                "name": person_name,
                                "auto_label": row[1],  # auto_label 列，供 KG 实体命名使用
                                "similarity": similarity,
                                "bbox": face_data.get(
                                    "bbox", []
                                ),  # 人脸框坐标 [x1, y1, x2, y2]
                                "confidence": face_data.get("confidence", 0.0),
                            }
                        )

                # 不在此处 commit，等文件复制和照片/人脸记录写入后一起提交

                # 4. Copy photo to storage
                workspace = get_workspace_path()

                # 构建存储路径
                relative_dir = build_photo_storage_path(category, source.name)
                target_dir = workspace / relative_dir
                target_dir.mkdir(parents=True, exist_ok=True)

                # 构建文件名（日期_时间）
                new_file_name = build_photo_file_name(source.name, exif.get("taken_at"))

                # 检查重名
                target_path = target_dir / new_file_name
                final_path = handle_photo_conflict(target_path)

                # 文件操作
                if mode == "move":
                    shutil.move(str(source), final_path)
                    logger.info(f"[INGEST_PHOTO] Moved to: {final_path}")
                elif mode == "reference":
                    final_path = str(source)
                    logger.info(f"[INGEST_PHOTO] Referenced: {final_path}")
                else:
                    shutil.copy2(str(source), final_path)
                    logger.info(f"[INGEST_PHOTO] Copied to: {final_path}")

                # 5. Create photo record
                photo_id = str(uuid.uuid4())

                # Generate L0 abstract (person_names 用于摘要，不用于文件名)
                person_names = [p["name"] for p in detected_persons]
                abstract = generate_l0_abstract(person_names, exif.get("taken_at"))

                conn.execute(
                    """INSERT INTO photos (id, file_path, taken_at, location, camera, abstract, ingested_at)
                       VALUES (?, ?, ?, ?, ?, ?, ?)""",
                    (
                        photo_id,
                        str(Path(final_path).resolve()),
                        exif.get("taken_at"),
                        exif.get("location"),
                        exif.get("camera"),
                        abstract,
                        now,
                    ),
                )

                # 6. Create face records
                if len(faces) != len(detected_persons):
                    logger.error(f"[INGEST_PHOTO] Face/person count mismatch: {len(faces)} vs {len(detected_persons)}")
                for i, (face_data, person) in enumerate(zip(faces, detected_persons)):
                    face_id = str(uuid.uuid4())
                    bbox_str = json.dumps(face_data["bbox"])
                    conn.execute(
                        """INSERT INTO faces (id, photo_id, person_id, embedding, bounding_box, confidence)
                           VALUES (?, ?, ?, ?, ?, ?)""",
                        (
                            face_id,
                            photo_id,
                            person["id"],
                            face_data["embedding"].tobytes(),
                            bbox_str,
                            face_data["confidence"],
                        ),
                    )

                # 7. Update co-occurrence relations（去重，防止同一人物多张人脸导致自共现）
                unique_persons = []
                seen_pids = set()
                for p in detected_persons:
                    if p["id"] not in seen_pids:
                        seen_pids.add(p["id"])
                        unique_persons.append(p)
                update_co_occurrences(unique_persons, exif.get("taken_at"))

                conn.commit()
        except Exception:
            if conn is not None:
                try:
                    conn.rollback()
                except Exception:
                    pass
            raise

        # 8. 同步到知识图谱（失败不影响照片入库）
        final_path_resolved = str(Path(final_path).resolve())
        kg_result = sync_photo_to_kg(final_path_resolved, abstract, detected_persons)
        kg_entities = kg_result.get("kg_entities", []) if isinstance(kg_result, dict) else []

        # 9. Unload face model to release memory (optional, for single photo)
        # For batch processing, keep model loaded
        # unload_face_model()

        logger.info(
            f"[INGEST_PHOTO] Completed: photo_id={photo_id}, persons={len(detected_persons)}"
        )

        return {
            "status": "success",
            "photo_id": photo_id,
            "file_path": str(Path(final_path).resolve()),
            "original_path": str(source),
            "category": category,
            "detected_persons": detected_persons,
            "abstract": abstract,
            "exif": exif,
            "lightrag_sync": kg_result,
            "kg_entities": kg_entities,
        }

    except Exception as e:
        logger.exception(f"[INGEST_PHOTO] Failed: {e}")
        # Rollback any uncommitted changes on the shared connection
        if conn is not None:
            try:
                conn.rollback()
            except Exception:
                pass
        # Clean up orphaned file if copy succeeded but DB failed
        if final_path is not None:
            try:
                if os.path.exists(final_path):
                    os.remove(final_path)
            except Exception:
                pass
        return {
            "status": "error",
            "error_code": "UNKNOWN_ERROR",
            "message": str(e),
            "kg_entities": [],
        }


def _merge_duplicate_person_entities(registry, target_name: str) -> None:
    """查询 KG 中是否存在与 target_name 同名的独立 person 实体，如有则合并。

    解决 H2 问题：ainsert 可能已创建独立的同名人名实体（如"任飞"），
    name_person 的 merge_entities 只合并旧 person:{uuid}，遗漏这些独立实体。

    流程：用 lightrag_query_data 的 keywords 参数搜索同名 person 实体，
    如发现多个同名 person 实体或名称近似匹配的 person 实体，执行额外合并。

    Args:
        registry: ToolRegistry 实例，用于调用 lightrag 工具。
        target_name: 目标人名实体名称。
    """
    try:
        query_fn = registry.get("lightrag-server/lightrag_query_data")
        merge_fn = registry.get("lightrag-server/lightrag_merge_entities")
        if not query_fn or not merge_fn:
            logger.debug("[NAME_PERSON] query_data or merge_entities not available, skip duplicate merge")
            return

        # 用 keywords 精确搜索，避免 LLM 提取延迟
        result = query_fn(query=target_name, mode="local", keywords=[target_name], top_k=20)

        if not result or result.get("status") in ("no_results", "error"):
            return

        # 提取实体列表
        data = result.get("data", result) if isinstance(result, dict) else {}
        if isinstance(data, list):
            entities = data
        else:
            entities = data.get("entities", []) if isinstance(data, dict) else []

        # 统计同名 person 实体出现次数和收集近似名称
        same_name_count = 0
        similar_name_entities: list[str] = []

        for entity in entities:
            if not isinstance(entity, dict):
                continue
            entity_name = entity.get("entity_name", "")
            entity_type = (entity.get("entity_type", "") or "").lower()

            if entity_type != "person":
                continue

            if entity_name == target_name:
                same_name_count += 1
            elif entity_name.startswith(target_name) and len(entity_name) - len(target_name) <= 5:
                # Conservative match: e.g. "任飞(人物)" starts with "任飞" and suffix ≤5 chars
                # Avoids false matches like "飞天" matching when target is "飞"
                similar_name_entities.append(entity_name)

        # 情况1：同名实体 > 1 个，说明 KG 中存在重复节点
        # （正常合并后应只有 1 个 target_name 实体）
        if same_name_count > 1:
            # 由于 LightRAG NetworkX 中同名节点唯一，这种情况理论上
            # 不应发生，但 VDB 层可能有重复向量记录
            logger.info(
                f"[NAME_PERSON] Found {same_name_count} duplicate person entities "
                f"named '{target_name}' in KG. This may indicate VDB inconsistency."
            )

        # 情况2：存在名称近似的 person 实体（ainsert LLM 提取的变体）
        # 执行合并：将近似实体合并到目标实体
        if similar_name_entities:
            # 去重
            unique_similar = list(dict.fromkeys(similar_name_entities))
            merge_fn(
                source_entities=unique_similar,
                target_entity=target_name,
            )
            logger.info(
                f"[NAME_PERSON] Merged {len(unique_similar)} similar person entities "
                f"({unique_similar}) → '{target_name}'"
            )
    except Exception as e:
        logger.warning(f"[NAME_PERSON] Duplicate person entity merge failed: {e}")


def _refresh_photo_abstracts_for_person(
    person_id: str, new_name: str, conn: sqlite3.Connection
) -> None:
    """人物改名后，更新相关照片的 abstract。

    只更新 DB 中照片的 abstract 字段，不再重新注入照片到 KG。
    避免 ainsert 重复产生新实体；人物实体改名已由 merge_entities 完成。

    流程：
    1. 从 DB 查找该人物出现的所有照片
    2. 重新生成 abstract（更新人名）
    3. 更新 DB 中的 abstract
    """
    # Find all photos that have faces from this person
    photo_rows = conn.execute(
        "SELECT DISTINCT f.photo_id, p.file_path, p.taken_at "
        "FROM faces f JOIN photos p ON f.photo_id = p.id "
        "WHERE f.person_id = ?",
        (person_id,),
    ).fetchall()

    if not photo_rows:
        return

    # Collect all DB updates first, then apply in one batch
    photo_updates: list[tuple[str, str, str, str]] = []  # (photo_id, file_path, new_abstract, taken_at)

    for photo_id, file_path, taken_at in photo_rows:
        # Get all person info for this photo (names are already updated in DB)
        face_persons = conn.execute(
            "SELECT p.id, p.name, p.auto_label FROM faces f "
            "JOIN persons p ON f.person_id = p.id "
            "WHERE f.photo_id = ?",
            (photo_id,),
        ).fetchall()

        # Build person_names list for abstract generation (deduplicate by person_id)
        person_names = []
        seen_person_ids: set[str] = set()
        for pid, pname, auto_label in face_persons:
            if pid in seen_person_ids:
                continue
            seen_person_ids.add(pid)
            is_unnamed = not pname or pname.startswith("未命名人物") or pname == auto_label
            display_name = auto_label if is_unnamed else pname
            if display_name:
                person_names.append(display_name)

        # Regenerate abstract with updated person names
        new_abstract = generate_l0_abstract(person_names, taken_at)

        photo_updates.append((photo_id, file_path, new_abstract, taken_at))

    # Batch all DB updates in one lock + one commit
    with _db_write_lock:
        for photo_id, _, new_abstract, _ in photo_updates:
            conn.execute(
                "UPDATE photos SET abstract = ? WHERE id = ?",
                (new_abstract, photo_id),
            )
        conn.commit()

    logger.info(
        f"[M7] Refreshed abstracts for {len(photo_updates)} photos after "
        f"naming person {person_id} as '{new_name}'"
    )


def name_person(person_id: str, name: str) -> dict:
    """Name an existing person."""
    try:
        conn = get_connection()

        # All DB reads/writes inside _db_write_lock to prevent TOCTOU
        with _db_write_lock:
            cursor = conn.execute(
                "SELECT id, name, auto_label, photo_count FROM persons WHERE id = ?", (person_id,)
            )
            row = cursor.fetchone()

            if not row:
                return {
                    "status": "error",
                    "error_code": "PERSON_NOT_FOUND",
                    "message": f"Person not found: {person_id}",
                }

            current_name = row[1]  # name column
            auto_label = row[2]    # auto_label column
            current_photo_count = row[3] or 0  # photo_count

            # Same-name detection: check if another person with the same name already exists
            if name and name != current_name:
                dup_cursor = conn.execute(
                    "SELECT id, name, auto_label, photo_count FROM persons WHERE name = ? AND id != ?",
                    (name, person_id),
                )
                dup_row = dup_cursor.fetchone()

                if dup_row:
                    return {
                        "status": "need_confirm",
                        "message": f"已存在名为\"{name}\"的人物",
                        "current_person": {
                            "person_id": person_id,
                            "name": current_name,
                            "auto_label": auto_label,
                            "photo_count": current_photo_count,
                        },
                        "existing_person": {
                            "person_id": dup_row[0],
                            "name": dup_row[1],
                            "auto_label": dup_row[2],
                            "photo_count": dup_row[3] or 0,
                        },
                        "merge_suggestion": {
                            "person_a_id": dup_row[0],   # existing_person — 保留（已命名、照片多）
                            "person_b_id": person_id,    # current_person — 合并后删除
                        },
                        "hint": "请确认：这是同一个人吗？如果是，请调用 merge_persons(person_a_id, person_b_id) 合并；如果只是同名，请换一个名字重新命名",
                    }

            conn.execute("UPDATE persons SET name = ? WHERE id = ?", (name, person_id))
            conn.commit()

        # KG: 通过 merge_entities 一步改名（旧实体删除+新实体创建+边迁移）
        kg_synced = False
        try:
            from agent.tool_registry import get_registry
            registry = get_registry()
            # 源实体：已命名用当前名，未命名用 auto_label
            source_entity = current_name if current_name else auto_label
            merge_fn = registry.get("lightrag-server/lightrag_merge_entities")
            inject_fn = registry.get("lightrag-server/lightrag_insert_custom_kg")
            # 先确保目标实体存在（merge_entities 不创建新实体，只迁移边）
            if inject_fn:
                inject_fn(
                    entities=[{
                        "entity_name": name,
                        "entity_type": "person",
                        "description": f"{name}，原名{source_entity}",
                    }],
                    relationships=[],
                    chunks=[],
                    source_id=f"rename_{source_entity}",
                )
            if merge_fn:
                merge_fn(
                    source_entities=[source_entity],
                    target_entity=name,
                )
                logger.info(f"[NAME_PERSON] KG renamed: {source_entity} → {name}")

                # H2 fix: 合并 KG 中已存在的同名独立实体（如 ainsert 创建的人名实体）
                _merge_duplicate_person_entities(registry, name)
                kg_synced = True
            else:
                logger.warning("[NAME_PERSON] lightrag_merge_entities not available in registry")
        except Exception as e:
            logger.warning(f"[NAME_PERSON] LightRAG rename failed: {e}")

        # M7: 改名后只更新 DB 中照片 abstract，不再重新注入照片到 KG
        # （避免 ainsert 重复产生新实体；人物实体改名已由 merge_entities 完成）
        try:
            _refresh_photo_abstracts_for_person(person_id, name, conn)
        except Exception as e:
            logger.warning(f"[NAME_PERSON] Photo abstract refresh failed: {e}")

        logger.info(f"[NAME_PERSON] Updated person {person_id} name to: {name}")

        return {
            "status": "success",
            "person_id": person_id,
            "name": name,
            "auto_label": auto_label,
            "kg_synced": kg_synced,
            "kg_rename": f"知识图谱实体名从「{source_entity}」改为「{name}」" if kg_synced else None,
        }

    except Exception as e:
        logger.exception(f"[NAME_PERSON] Failed: {e}")
        return {
            "status": "error",
            "error_code": "UNKNOWN_ERROR",
            "message": str(e),
            "kg_rename": None,
        }


def merge_persons(person_a_id: str, person_b_id: str) -> dict:
    """Merge two persons into one (keeping person_a's name)."""
    conn = None

    if person_a_id == person_b_id:
        return {
            "status": "error",
            "error_code": "SAME_PERSON_ID",
            "message": "Cannot merge a person with themselves",
            "merged_into": None,
            "deleted_person_id": None,
            "kg_rename": None,
        }

    try:
        conn = get_connection()

        # Get both persons (包含 threshold_adjustment)
        cursor = conn.execute(
            "SELECT id, name, auto_label, center_embedding, threshold_adjustment, photo_count FROM persons WHERE id IN (?, ?)",
            (person_a_id, person_b_id),
        )
        rows = cursor.fetchall()

        if len(rows) != 2:
            return {
                "status": "error",
                "error_code": "PERSON_NOT_FOUND",
                "message": "One or both persons not found",
                "kg_rename": None,
            }

        # Find person_a and person_b
        persons = {row[0]: row for row in rows}
        person_a = persons.get(person_a_id)
        person_b = persons.get(person_b_id)

        if not person_a or not person_b:
            return {
                "status": "error",
                "error_code": "PERSON_NOT_FOUND",
                "message": "One or both persons not found",
                "kg_rename": None,
            }

        # Keep person_a's name
        name_a = person_a[1]
        auto_label_a = person_a[2]

        # Merge embeddings
        embedding_a = (
            np.frombuffer(person_a[3], dtype=np.float32) if person_a[3] else None
        )
        embedding_b = (
            np.frombuffer(person_b[3], dtype=np.float32) if person_b[3] else None
        )

        if embedding_a is not None and embedding_b is not None:
            # Weighted average based on photo count
            count_a = person_a[5] or 1
            count_b = person_b[5] or 1
            total = count_a + count_b
            merged_embedding = (embedding_a * count_a + embedding_b * count_b) / total
        elif embedding_a is not None:
            merged_embedding = embedding_a
        elif embedding_b is not None:
            merged_embedding = embedding_b
        else:
            merged_embedding = None

        # Calculate threshold adjustment based on original similarity
        threshold_adjustment = 0.0
        if embedding_a is not None and embedding_b is not None:
            similarity = cosine_similarity(embedding_a, embedding_b)
            # If similarity was below threshold but user confirms same person,
            # adjust threshold to make future matches easier
            if similarity < 0.7:
                threshold_adjustment = (0.7 - similarity) + 0.05
                threshold_adjustment = min(threshold_adjustment, 0.3)  # Cap at 0.3

        # 写事务：合并人物数据并提交
        try:
            with _db_write_lock:
                # Update all faces from person_b to person_a
                conn.execute(
                    "UPDATE faces SET person_id = ? WHERE person_id = ?",
                    (person_a_id, person_b_id),
                )

                # Calculate merged_count from actual face records (handles overlapping photos)
                merged_count = conn.execute(
                    "SELECT COUNT(DISTINCT photo_id) FROM faces WHERE person_id = ?",
                    (person_a_id,),
                ).fetchone()[0]

                # Update person_a with merged data（始终执行，即使 embedding 为 None）
                conn.execute(
                    """UPDATE persons SET center_embedding = ?, photo_count = ?,
                       threshold_adjustment = ?, last_seen = ? WHERE id = ?""",
                    (
                        merged_embedding.tobytes() if merged_embedding is not None else None,
                        merged_count,
                        max(threshold_adjustment, person_a[4] if person_a[4] else 0),
                        datetime.now().isoformat(),
                        person_a_id,
                    ),
                )

                # Clean up co_occurrences referencing person_b
                # 1. Delete the pair (person_a, person_b) — self-co-occurrence after merge
                pair = tuple(sorted([person_a_id, person_b_id]))
                conn.execute(
                    "DELETE FROM co_occurrences WHERE person_a_id = ? AND person_b_id = ?",
                    pair,
                )
                # 2. Transfer person_b's co-occurrences to person_a (merge counts)
                #    For each row where person_b appears with a third person X,
                #    add the count to person_a's existing row with X (or create new row)
                b_rows = conn.execute(
                    "SELECT person_a_id, person_b_id, count, first_seen, last_seen FROM co_occurrences WHERE person_a_id = ? OR person_b_id = ?",
                    (person_b_id, person_b_id),
                ).fetchall()
                for row_a, row_b, count, first_seen, last_seen in b_rows:
                    # Determine the third person X
                    other_id = row_b if row_a == person_b_id else row_a
                    if other_id == person_a_id:
                        continue  # Skip self-pair (already deleted above)
                    co_pair = tuple(sorted([person_a_id, other_id]))
                    # Try to add count to existing row
                    existing = conn.execute(
                        "SELECT count, first_seen, last_seen FROM co_occurrences WHERE person_a_id = ? AND person_b_id = ?",
                        co_pair,
                    ).fetchone()
                    if existing:
                        existing_first = existing[1]
                        existing_last = existing[2]
                        merged_first = min(existing_first, first_seen) if existing_first and first_seen else (existing_first or first_seen)
                        merged_last = max(existing_last, last_seen) if existing_last and last_seen else (existing_last or last_seen)
                        conn.execute(
                            "UPDATE co_occurrences SET count = count + ?, first_seen = ?, last_seen = ? WHERE person_a_id = ? AND person_b_id = ?",
                            (count, merged_first, merged_last, co_pair[0], co_pair[1]),
                        )
                    else:
                        conn.execute(
                            "INSERT INTO co_occurrences (person_a_id, person_b_id, count, first_seen, last_seen) VALUES (?, ?, ?, ?, ?)",
                            (co_pair[0], co_pair[1], count, first_seen, last_seen),
                        )
                # 3. Delete all remaining rows referencing person_b
                conn.execute(
                    "DELETE FROM co_occurrences WHERE person_a_id = ? OR person_b_id = ?",
                    (person_b_id, person_b_id),
                )

                # Delete person_b
                conn.execute("DELETE FROM persons WHERE id = ?", (person_b_id,))

                conn.commit()
        except Exception:
            if conn is not None:
                try:
                    conn.rollback()
                except Exception:
                    pass
            raise

        # 同步 LightRAG：更新目标实体描述，合并源实体关系到目标实体
        kg_synced = False
        try:
            from agent.tool_registry import get_registry

            registry = get_registry()
            # KG 实体名：已命名用 name，未命名用 auto_label
            kg_name_a = name_a if name_a else auto_label_a
            name_b = person_b[1]
            auto_label_b = person_b[2]
            kg_name_b = name_b if name_b else auto_label_b

            # 1. 用 inject_custom_kg 确保目标实体存在并更新描述（chunks=[] → 不触发 LLM）
            inject_fn = registry.get("lightrag-server/lightrag_insert_custom_kg")
            if not inject_fn:
                logger.warning("[MERGE_PERSONS] lightrag_insert_custom_kg not available, skipping KG sync")
            else:
                inject_fn(
                    entities=[{
                        "entity_name": kg_name_a,
                        "entity_type": "person",
                        "description": f"{kg_name_a}，合并自{kg_name_b}",
                    }],
                    relationships=[],
                    chunks=[],
                    source_id=f"merge_{kg_name_a}",
                )

                # 2. 合并：kg_name_b 的边迁移到 kg_name_a（依赖步骤1确保目标实体存在）
                merge_fn = registry.get("lightrag-server/lightrag_merge_entities")
                if merge_fn:
                    merge_fn(
                        source_entities=[kg_name_b],
                        target_entity=kg_name_a,
                    )
                    kg_synced = True
                logger.info(f"[MERGE_PERSONS] Merged KG entity {kg_name_b} into {kg_name_a}")
                # M5: After merging, if both source entities had edges to the same
                # third-party entity, duplicate edges may exist. Dedup requires
                # modifying LightRAG core code, so we log a warning instead.
                logger.debug(
                    f"[MERGE_PERSONS] Note: if {kg_name_a} and {kg_name_b} both had "
                    f"edges to the same third-party entity, duplicate edges may now "
                    f"exist on {kg_name_a}. Dedup requires LightRAG core changes."
                )
        except Exception as e:
            logger.warning(f"[MERGE_PERSONS] LightRAG sync failed: {e}")

        logger.info(f"[MERGE_PERSONS] Merged {person_b_id} into {person_a_id}")

        return {
            "status": "success",
            "merged_into": person_a_id,
            "name": name_a if name_a else auto_label_a,
            "photo_count": merged_count,
            "deleted_person_id": person_b_id,
            "kg_rename": f"知识图谱实体名从「{kg_name_b}」改为「{kg_name_a}」" if kg_synced else None,
        }

    except Exception as e:
        logger.exception(f"[MERGE_PERSONS] Failed: {e}")
        if conn is not None:
            try:
                conn.rollback()
            except Exception:
                pass
        return {
            "status": "error",
            "error_code": "UNKNOWN_ERROR",
            "message": str(e),
            "kg_rename": None,
        }


# ============== 照片批量处理 ==============

PHOTO_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".heic", ".heif"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".avi", ".mkv", ".webm", ".flv", ".wmv", ".m4v"}


def is_photo(file_path: str) -> bool:
    """判断是否为照片文件"""
    return Path(file_path).suffix.lower() in PHOTO_EXTENSIONS


def is_video(file_path: str) -> bool:
    """判断是否为视频文件"""
    return Path(file_path).suffix.lower() in VIDEO_EXTENSIONS


DOCUMENT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".txt",
    ".md",
    ".xlsx",
    ".xls",
    ".pptx",
    ".ppt",
}


def is_document(file_path: str) -> bool:
    """判断是否为文档类文件"""
    return Path(file_path).suffix.lower() in DOCUMENT_EXTENSIONS


# LightRAG 知识图谱支持的文件扩展名
KG_SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".log",
    ".pdf", ".docx", ".pptx", ".xlsx",
    ".html", ".htm",
}

# OLE2 复合文档签名（旧版 .doc 格式）
OLE2_MAGIC = b"\xd0\xcf\x11\xe0"


def check_kg_supported(file_path: str) -> tuple[bool, str]:
    """检查文件格式是否支持知识图谱入库

    Returns:
        (supported, reason) 元组：
        - supported: True 表示支持，False 表示不支持
        - reason: 不支持时的原因说明
    """
    suffix = Path(file_path).suffix.lower()

    # 无扩展名
    if not suffix:
        return (False, "无扩展名的文件不支持知识图谱入库")

    # 扩展名不在白名单中
    if suffix not in KG_SUPPORTED_EXTENSIONS:
        return (False, f"{suffix} 格式不支持知识图谱入库")

    # 检测假 .docx（WPS 创建的 OLE2 格式文件，扩展名为 .docx 但实际是旧版 .doc）
    # OLE2 签名前 4 字节: D0 CF 11 E0（完整签名 8 字节，4 字节足以区分 ZIP vs OLE2）
    if suffix == ".docx":
        try:
            with open(file_path, "rb") as f:
                header = f.read(4)
            if header == OLE2_MAGIC:
                return (False, "该 .docx 文件实际为旧版 .doc 格式（WPS 创建），不支持知识图谱入库")
        except Exception as e:
            logger.warning(f"[KG_CHECK] 无法读取文件头检测 OLE2: {e}")

    return (True, "")


def calculate_file_hash(file_path: str) -> str:
    """计算文件哈希"""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


def read_file_content(path: str, max_chars: int = 20000) -> str:
    """读取文件内容，支持多种格式，超过max_chars截断"""
    suffix = Path(path).suffix.lower()
    text = ""

    if suffix in {".txt", ".md", ".csv", ".json", ".log"}:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read(max_chars)

    elif suffix == ".pdf":
        try:
            from pypdf import PdfReader

            with open(path, "rb") as f:
                reader = PdfReader(f)
                parts = []
                total = 0
                for page in reader.pages:
                    t = page.extract_text() or ""
                    parts.append(t)
                    total += len(t)
                    if total >= max_chars:
                        break
                text = " ".join(parts)
        except Exception as e:
            logger.warning(f"[READ] PDF读取失败: {e}")

    elif suffix == ".docx":
        try:
            from docx import Document
            doc = Document(path)
            parts = []
            total = 0
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
                    total += len(para.text)
                    if total >= max_chars:
                        break
            text = "\n".join(parts)
        except Exception as e:
            logger.warning(f"[READ] DOCX读取失败: {e}")

    elif suffix in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            total = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    t = getattr(shape, "text", "")
                    t_stripped = str(t).strip()
                    if t_stripped:
                        parts.append(t_stripped)
                        total += len(t_stripped)
                        if total >= max_chars:
                            break
                if total >= max_chars:
                    break
            text = "\n".join(parts)
        except Exception as e:
            logger.warning(f"[READ] PPTX读取失败: {e}")

    elif suffix == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(path, data_only=True)
            parts = []
            total = 0
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(c) for c in row if c is not None])
                    if row_text.strip():
                        parts.append(row_text)
                        total += len(row_text)
                        if total >= max_chars:
                            break
                if total >= max_chars:
                    break
            text = "\n".join(parts)
        except Exception as e:
            logger.warning(f"[READ] XLSX读取失败: {e}")

    elif suffix in {".html", ".htm"}:
        try:
            from html.parser import HTMLParser

            class _TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self._texts = []

                def handle_data(self, data):
                    if data.strip():
                        self._texts.append(data.strip())

            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                html = f.read()
            parser = _TextExtractor()
            parser.feed(html)
            text = "\n".join(parser._texts)
        except Exception as e:
            logger.warning(f"[READ] HTML读取失败: {e}")

    else:
        # Unknown format — check if binary before attempting text read
        try:
            with open(path, "rb") as f:
                chunk = f.read(1024)
            if b"\x00" in chunk:
                return ""  # binary file, skip text extraction
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read(max_chars)
        except Exception:
            return ""

    return text[:max_chars]


def calculate_content_similarity(file1: str, file2: str) -> float:
    """
    计算文档内容相似度（语义向量 + 余弦相似度）
    使用共享向量服务
    """
    try:
        logger.info(f"[SIMILARITY] 开始计算: {file1} vs {file2}")

        # 读取文件内容
        text1 = read_file_content(file1)
        text2 = read_file_content(file2)
        logger.info(f"[SIMILARITY] 文本长度: {len(text1)} vs {len(text2)}")

        if not text1 or not text2:
            logger.info("[SIMILARITY] 空文本，返回 0.0")
            return 0.0

        # 截断过长的文本（避免内存问题）
        max_len = 8000
        if len(text1) > max_len:
            text1 = text1[:max_len]
        if len(text2) > max_len:
            text2 = text2[:max_len]

        # 调用共享向量服务计算相似度
        logger.info("[SIMILARITY] 调用共享向量服务...")
        result = call_embedding_service("/similarity", {"text1": text1, "text2": text2})

        if result and "similarity" in result:
            similarity = result["similarity"]
            logger.info(f"[SIMILARITY] 结果: {similarity:.4f}")
            return similarity
        else:
            logger.warning("[SIMILARITY] 向量服务不可用，返回 0.0")
            return 0.0

    except Exception as e:
        logger.exception(f"[SIMILARITY] 失败: {e}")
        return 0.0


def build_storage_path(
    category: str, file_name: str, file_type: str = "documents"
) -> str:
    """构建存储路径（相对路径）"""
    prefs = get_preferences()
    structure = prefs["storage"]["structure"].get(file_type, "{year}/{category}")

    now = datetime.now()
    path = structure.replace("{year}", str(now.year))
    path = path.replace("{month}", f"{now.month:02d}")
    path = path.replace("{date}", now.strftime("%Y-%m-%d"))
    path = path.replace("{category}", category)

    return path


def handle_conflict(target_path: Path, source_path: Path) -> tuple[str, str]:
    """
    处理文件冲突
    返回: (最终路径, 动作说明)
    """
    logger.info(f"[CONFLICT] 检查: {target_path}")
    if not target_path.exists():
        logger.info("[CONFLICT] 文件不存在，创建新文件")
        return str(target_path), "created"

    logger.info("[CONFLICT] 文件已存在，读取冲突配置...")
    prefs = get_preferences()
    conflict_config = prefs["storage"]["conflict"]

    if is_document(str(target_path)):
        logger.info("[CONFLICT] 文档类型，计算哈希...")
        hash_existing = calculate_file_hash(str(target_path))
        hash_new = calculate_file_hash(str(source_path))
        logger.info(
            f"[CONFLICT] 哈希: existing={hash_existing[:8]}... new={hash_new[:8]}..."
        )

        if hash_existing == hash_new:
            logger.info("[CONFLICT] 哈希相同，跳过")
            return str(target_path), "skipped"

        # 哈希不同，计算语义相似度
        logger.info("[CONFLICT] 哈希不同，计算语义相似度...")
        threshold = conflict_config["document"]["similarity_threshold"]
        similarity = calculate_content_similarity(str(target_path), str(source_path))
        logger.info(f"[CONFLICT] 相似度: {similarity:.4f}, 阈值: {threshold}")

        if similarity > threshold:
            logger.info("[CONFLICT] 相似，版本管理")
            return create_version(target_path), "versioned"
        else:
            logger.info("[CONFLICT] 不相似，改名")
            return rename_file(target_path), "renamed"
    else:
        logger.info("[CONFLICT] 非文档类型，改名")
        return rename_file(target_path), "renamed"


def create_version(file_path: Path) -> str:
    """创建版本：原文件改名，返回原文件名路径"""
    timestamp = datetime.now().strftime("%Y%m%d")
    stem = file_path.stem
    suffix = file_path.suffix

    version_name = f"{stem}_v1_{timestamp}{suffix}"
    version_path = file_path.parent / version_name

    index = 1
    while version_path.exists():
        index += 1
        version_name = f"{stem}_v{index}_{timestamp}{suffix}"
        version_path = file_path.parent / version_name

    shutil.move(str(file_path), str(version_path))
    logger.info(f"版本管理: {file_path} -> {version_path}")

    return str(file_path)


def rename_file(file_path: Path) -> str:
    """改名：返回新的文件名"""
    stem = file_path.stem
    suffix = file_path.suffix
    parent = file_path.parent

    index = 1
    new_path = parent / f"{stem}_{index}{suffix}"
    while new_path.exists():
        index += 1
        new_path = parent / f"{stem}_{index}{suffix}"

    logger.info(f"文件改名预留: {file_path} -> {new_path}")
    return str(new_path)


# ============== 有状态目录入库 ==============


def _normalize_session_key(path: str) -> str:
    """归一化路径作为会话 key（统一分隔符 + 小写）"""
    return str(Path(path).resolve()).replace("\\", "/").lower()


def _scan_directory(source: Path) -> list[dict]:
    """扫描目录，分类文件为 image/document/skipped
    
    Returns:
        [{"path": str, "type": "image"|"document"|"skipped", "reason": str}, ...]
    """
    files = []
    seen_paths = set()
    
    # 收集所有文件
    for item in source.rglob("*"):
        if not item.is_file():
            continue
        key = str(item.resolve()).lower()
        if key in seen_paths:
            continue
        seen_paths.add(key)
        
        suffix = item.suffix.lower()
        if suffix in PHOTO_EXTENSIONS:
            files.append({"path": str(item), "type": "image"})
        elif suffix in DOCUMENT_EXTENSIONS:
            files.append({"path": str(item), "type": "document"})
        elif suffix in VIDEO_EXTENSIONS:
            files.append({"path": str(item), "type": "skipped", "reason": "视频文件暂不支持入库"})
        else:
            files.append({"path": str(item), "type": "skipped", "reason": f"不支持的文件格式: {suffix}"})
    
    return files


def _find_next_processable(all_files: list[dict], offset: int) -> int | None:
    """从 offset 开始找下一个可处理文件（image 或 document），返回索引或 None"""
    for i in range(offset, len(all_files)):
        if all_files[i]["type"] in ("image", "document"):
            return i
    return None


def _build_success_summary(all_files: list[dict], processed: list[dict]) -> dict:
    """构建最终成功摘要"""
    photos = [p for p in processed if p["type"] == "image"]
    documents = [p for p in processed if p["type"] == "document"]
    skipped = [f for f in all_files if f["type"] == "skipped"]
    errors = [p for p in processed if isinstance(p.get("result"), dict) and p["result"].get("status") == "error"]

    return {
        "status": "success",
        "total": len(all_files),
        "photos": len(photos),
        "documents": len(documents),
        "skipped": len(skipped),
        "errors": len(errors),
        "details": {
            "photos": [{"file": Path(p["file"]).name, "result": p["result"]} for p in photos],
            "documents": [{"file": Path(p["file"]).name, "result": p["result"]} for p in documents],
            "skipped": [{"file": Path(f["path"]).name, "reason": f.get("reason", "")} for f in skipped],
            "errors": [{"file": Path(p["file"]).name, "message": p["result"].get("message", "")} for p in errors],
        },
    }


def _process_next_file(session: dict, category: str = "") -> dict:
    """处理 session 中的下一个文件，更新进度，返回结果

    自动循环逻辑：
    - 照片：始终自动循环（不需要分类）
    - 文档 + session["auto_category"] 非空：用 auto_category 自动循环
    - 文档 + category 非空（子Agent刚回答的）：只处理当前文档，不传播到后续
    - 文档 + 都为空：返回 need_category 等待交互

    Args:
        session: 会话状态 dict（含 auto_category 字段）
        category: 子Agent回答的当前文档分类（只处理当前文档，不传播）

    Returns:
        处理结果 dict（need_category / success / error）
    """
    all_files = session["all_files"]
    mode = session["mode"]
    auto_category = session.get("auto_category", "")

    # 循环处理文件，直到遇到 need_category 或全部完成
    while True:
        offset = session["offset"]
        next_idx = _find_next_processable(all_files, offset)
        if next_idx is None:
            # 全部处理完毕
            return _build_success_summary(all_files, session["processed"])

        current = all_files[next_idx]

        try:
            if current["type"] == "image":
                # 照片不需要分类判断，始终自动继续
                session["offset"] = next_idx + 1  # 先推进 offset，防止 append 失败导致重复处理
                result = ingest_photo(current["path"], category=None, mode=mode)
                session["processed"].append({"file": current["path"], "type": "image", "result": result})
                # 继续循环处理下一个文件
                continue

            elif current["type"] == "document":
                # 优先使用 auto_category（初始化时指定的全量分类）
                if auto_category:
                    # 有全量分类，自动入库并继续循环
                    session["offset"] = next_idx + 1  # 先推进 offset，防止 append 失败导致重复处理
                    result = ingest_document(file_path=current["path"], category=auto_category, mode=mode)
                    session["processed"].append({"file": current["path"], "type": "document", "result": result})
                    # 继续循环处理下一个文件
                    continue
                elif category:
                    # 子Agent回答的分类，只处理当前文档，不传播到后续
                    session["offset"] = next_idx + 1  # 先推进 offset，防止 append 失败导致重复处理
                    result = ingest_document(file_path=current["path"], category=category, mode=mode)
                    session["processed"].append({"file": current["path"], "type": "document", "result": result})
                    # 清空 category，后续文档根据 auto_category 判断（为空则返回 need_category）
                    category = ""
                    # 继续循环，下一个文档会根据 auto_category 决定行为
                    continue
                else:
                    # 需要分类 — 先读取内容预览
                    doc_result = ingest_document(file_path=current["path"], category="", mode=mode)

                    # 如果 ingest_document 返回 error（文件无法处理），跳过此文件继续下一个
                    if doc_result.get("status") == "error":
                        session["processed"].append({"file": current["path"], "type": "document", "result": doc_result})
                        session["offset"] = next_idx + 1
                        # 继续循环处理下一个文件
                        continue

                    # 从 ingest_document 返回结果中提取预览（已包含 20000 字符）
                    preview = doc_result.get("preview", "")
                    if not preview:
                        # ingest_document 未能读取内容，记录错误并跳过
                        logger.warning(f"文档预览为空: {current['path']}")
                        session["processed"].append({
                            "file": current["path"],
                            "type": "document",
                            "result": {"status": "error", "message": "无法读取文档预览: 内容为空"}
                        })
                        session["offset"] = next_idx + 1
                        continue
                    available_categories = doc_result.get("available_categories", ["其他"])

                    return {
                        "status": "need_category",
                        "total": len(all_files),
                        "current_file": Path(current["path"]).name,
                        "current_file_path": current["path"],
                        "preview": preview,
                        "available_categories": available_categories,
                        "message": f"请从 available_categories 中选择分类，然后再次调用 ingest(path=\"{current['path']}\", category=分类名) 继续",
                    }

            else:
                # 未知类型，跳过
                session["offset"] = next_idx + 1
                continue

        except Exception as e:
            logger.error(f"处理文件异常: {current['path']}, 错误: {e}")
            session["processed"].append({
                "file": current["path"],
                "type": current.get("type", "unknown"),
                "result": {"status": "error", "message": str(e)}
            })
            session["offset"] = next_idx + 1
            continue


def ingest(path: str, category: str = "", mode: str = "copy", action: str = "") -> dict:
    """有状态统一入库工具
    
    单文件：无状态，直接入库
    目录：有状态三阶段交互（start → interact → ... → success/abort）
    """
    source = Path(os.path.expanduser(path))
    
    if not source.exists():
        return {
            "status": "error",
            "error_code": "PATH_NOT_FOUND",
            "message": f"路径不存在: {path}",
        }
    
    # 单文件：无状态直接入库
    if source.is_file():
        if is_photo(str(source)):
            return ingest_photo(str(source), category=category or None, mode=mode)
        else:
            return ingest_document(file_path=str(source), category=category, mode=mode)
    
    # 目录：有状态处理
    if not source.is_dir():
        return {
            "status": "error",
            "error_code": "INVALID_PATH",
            "message": f"无效路径: {path}",
        }
    
    session_key = _normalize_session_key(path)
    
    # action="abort": 中止会话
    if action == "abort":
        session = _ingest_sessions.pop(session_key, None)
        if session is None:
            return {
                "status": "error",
                "message": "没有活跃的入库会话",
            }
        processed_count = len(session["processed"])
        return {
            "status": "aborted",
            "message": f"入库已中止，已处理 {processed_count} 个文件",
            "processed_count": processed_count,
        }
    
    # action="start": 初始化会话
    if action == "start":
        # 如果已有会话，先清理
        if session_key in _ingest_sessions:
            logger.warning(f"[INGEST] 覆盖已有会话: {session_key}")

        all_files = _scan_directory(source)
        if not all_files:
            return {
                "status": "error",
                "error_code": "NO_FILES_FOUND",
                "message": f"目录中没有找到可处理的文件: {path}",
            }

        # 仅当目录中包含文档文件时才校验 category（纯照片目录不需要分类）
        has_documents = any(f["type"] == "document" for f in all_files)
        if category and has_documents:
            prefs = get_preferences()
            available_categories = prefs.get("categories", {}).get("documents", ["其他"])
            if category not in available_categories:
                return {
                    "status": "error",
                    "error_code": "INVALID_CATEGORY",
                    "message": f"分类 '{category}' 不在可选列表中，可选分类: {', '.join(available_categories)}",
                }
        
        image_count = sum(1 for f in all_files if f["type"] == "image")
        doc_count = sum(1 for f in all_files if f["type"] == "document")
        skip_count = sum(1 for f in all_files if f["type"] == "skipped")
        
        session = {
            "all_files": all_files,
            "offset": 0,
            "processed": [],
            "mode": mode,
            "auto_category": category,  # 初始化时的分类，非空表示全量自动循环
        }
        _ingest_sessions[session_key] = session
        
        # 自动处理文件（category 非空时自动循环处理完所有文件）
        result = _process_next_file(session, category=category)

        # 如果全部处理完毕，清理会话
        if result["status"] == "success":
            _ingest_sessions.pop(session_key, None)
            return result

        # need_category 时附加概览信息
        if result["status"] == "need_category":
            result["message"] = f"发现 {image_count} 张图片、{doc_count} 个文档、{skip_count} 个跳过文件\n" + result["message"]

        return result
    
    # action="" 或 "interact": 继续交互
    if action in ("", "interact"):
        session = _ingest_sessions.get(session_key)
        if session is None:
            return {
                "status": "error",
                "message": "会话未初始化，请先调用 ingest(path, action='start')",
            }
        
        # 如果传了 category 且当前处于 need_category 状态
        # （offset 指向的文件就是需要分类的那个文档）
        if category:
            # 找到当前需要分类的文件
            next_idx = _find_next_processable(session["all_files"], session["offset"])
            if next_idx is not None and session["all_files"][next_idx]["type"] == "document":
                current = session["all_files"][next_idx]
                # 验证 category 是否在可选列表中
                prefs = get_preferences()
                available_categories = prefs.get("categories", {}).get("documents", ["其他"])
                if category not in available_categories:
                    # 重新读取预览
                    preview = ""
                    try:
                        preview = read_file_content(current["path"])[:20000]
                    except Exception:
                        pass
                    return {
                        "status": "need_category",
                        "total": len(session["all_files"]),
                        "current_file": Path(current["path"]).name,
                        "current_file_path": current["path"],
                        "preview": preview,
                        "available_categories": available_categories,
                        "message": f"分类 '{category}' 不在可选列表中，请从 available_categories 中选择",
                    }
                
                # 有分类，处理当前文档
                result = ingest_document(file_path=current["path"], category=category, mode=session["mode"])
                session["processed"].append({"file": current["path"], "type": "document", "result": result})
                session["offset"] = next_idx + 1

                # 继续处理（不传 category，让函数根据 session["auto_category"] 决定后续行为）
                next_result = _process_next_file(session)
                if next_result["status"] == "success":
                    _ingest_sessions.pop(session_key, None)
                return next_result

        # 无 category，继续处理下一个文件
        result = _process_next_file(session)
        if result["status"] == "success":
            _ingest_sessions.pop(session_key, None)
        return result
    
    return {
        "status": "error",
        "message": f"未知的 action: {action}",
    }


# ============== 工具实现 ==============


def ingest_document(file_path: str, category: str = "", mode: str = "copy") -> dict:
    """文档入库工具 — 文件搬运 + 知识图谱入库（格式不支持时跳过KG）

    不传 category 时，读取文件内容返回 need_category 状态+内容预览，
    由调用方判断分类后再次调用传入 category 完成入库。

    自动检测路径类型（目录/照片/文档）：
    - 目录：检查是否包含照片，转到照片批量处理
    - 照片：转到照片入库流程
    - 文档：文件搬运 + 检查格式支持后提交LightRAG（不支持则跳过KG）

    返回 status:
    - need_category: 需要调用方判断分类（附带内容预览和 kg_supported 标记）
    - success: 入库完成（lightrag: inserted | unsupported | error | skipped）
    - error: 失败
    """
    try:
        logger.info(f"[INGEST] 开始处理: {file_path}")
        source = Path(os.path.expanduser(file_path))
        if not source.exists():
            logger.error(f"[INGEST] 文件不存在: {file_path}")
            return {
                "status": "error",
                "error_code": "FILE_NOT_FOUND",
                "message": f"文件不存在: {file_path}",
                "suggestion": "请检查文件路径是否正确",
                "kg_entities": [],
            }

        # 目录应使用 ingest 工具处理
        if source.is_dir():
            return {
                "status": "error",
                "error_code": "IS_DIRECTORY",
                "message": "请使用 ingest 工具处理目录，ingest_document 仅处理单文件",
                "kg_entities": [],
            }

        if is_photo(str(source)):
            logger.info("[INGEST] 检测到照片文件，转到照片入库")
            from niu_photo_server import ingest_photo
            return ingest_photo(str(source), category=category or None, mode=mode)

        # ---- 文档文件处理 ----

        # 检查格式是否支持知识图谱入库（提前告知子 Agent）
        kg_ok, kg_reason = check_kg_supported(file_path)

        # No category → read file content and ask caller to classify
        if not category:
            prefs = get_preferences()
            available_categories = prefs.get("categories", {}).get("documents", ["其他"])
            content = read_file_content(file_path)
            kg_note = f"\n注意：{kg_reason}" if not kg_ok else ""
            if content:
                preview = content[:20000] if len(content) > 20000 else content
                return {
                    "status": "need_category",
                    "message": f"请根据以下内容判断文档分类目录，然后再次调用 ingest 工具并传入 category 参数。对于目录入库使用相同的 path 参数继续会话；对于单文件入库使用相同的 file_path 参数。\n\n文件: {file_path}\n内容预览:\n{preview}\n可选分类: {', '.join(available_categories)}{kg_note}",
                    "file_path": file_path,
                    "mode": mode,
                    "content_length": len(content),
                    "available_categories": available_categories,
                    "kg_supported": kg_ok,
                    "preview": preview,
                }
            else:
                ext = source.suffix.lower()
                size = source.stat().st_size
                return {
                    "status": "need_category",
                    "message": f"无法读取文件内容，请根据文件信息判断分类目录，然后再次调用 ingest 工具并传入 category 参数。对于目录入库使用相同的 path 参数继续会话；对于单文件入库使用相同的 file_path 参数。\n\n文件: {file_path}\n格式: {ext}\n大小: {size} 字节\n可选分类: {', '.join(available_categories)}{kg_note}",
                    "file_path": file_path,
                    "mode": mode,
                    "available_categories": available_categories,
                    "kg_supported": kg_ok,
                }
        logger.info("[INGEST] 读取配置...")
        workspace = get_workspace_path()

        logger.info("[INGEST] 构建存储路径...")
        relative_dir = build_storage_path(category, source.name, "documents")
        target_dir = workspace / relative_dir
        target_dir.mkdir(parents=True, exist_ok=True)

        target_path = target_dir / source.name
        logger.info(f"[INGEST] 目标路径: {target_path}")

        logger.info("[INGEST] 检查冲突...")
        final_path, action = handle_conflict(target_path, source)
        logger.info(f"[INGEST] 冲突处理结果: {action}")

        if action == "skipped":
            logger.info("[INGEST] 文件已存在，检查 LightRAG 写入状态...")
            content_length = 0
            lightrag_result = None

            # 使用前面已检查的 KG 格式支持结果
            if not kg_ok:
                lr_status = "unsupported"
                lr_msg = kg_reason
                logger.info(f"[INGEST] KG unsupported (skipped file): {kg_reason}")
                content_length = Path(final_path).stat().st_size
            else:
                lr_status = "skipped"
                lr_msg = ""
                try:
                    from agent.tool_registry import get_registry

                    registry = get_registry()
                    insert_file_tool = registry.get("lightrag-server/lightrag_insert_file")
                    if insert_file_tool:
                        lightrag_result = insert_file_tool(
                            file_path=str(Path(final_path).resolve()),
                            doc_id=str(Path(final_path).resolve()),
                        )
                        lr_status = lightrag_result.get("status", "unknown")
                        logger.info(f"[INGEST] LightRAG insert_file (skipped path): {lr_status}")
                        content_length = Path(final_path).stat().st_size
                except Exception as lr_err:
                    logger.warning(f"[INGEST] LightRAG insert_file 失败（不影响文件入库）: {lr_err}")

            # 确定 lightrag 状态
            if lr_status == "unsupported":
                lr_final = "unsupported"
                lr_final_msg = lr_msg
            elif lightrag_result and lightrag_result.get("status") in ("success", "ok"):
                lr_final = "inserted"
                lr_final_msg = ""
            elif lightrag_result and lightrag_result.get("status") == "error":
                lr_final = "error"
                lr_final_msg = lightrag_result.get("message", "")
            else:
                lr_final = "skipped"
                lr_final_msg = ""

            # 根据知识图谱状态生成提示
            if lr_final == "unsupported":
                note = f"文件已存在，但{lr_final_msg}"
            elif lr_final == "error":
                note = f"文件已存在，知识图谱入库失败：{lr_final_msg}"
            else:
                note = "文件已存在，已补全 LightRAG 写入"

            return {
                "status": "success",
                "action": "skipped",
                "file_path": str(Path(final_path).resolve()),
                "original_path": str(source),
                "category": category,
                "content_length": content_length,
                "lightrag": lr_final,
                "lightrag_message": lr_final_msg or None,
                "note": note,
                "kg_entities": [],
            }

        logger.info(f"[INGEST] 执行文件操作: {mode}")
        if mode == "copy":
            shutil.copy2(str(source), final_path)
        elif mode == "move":
            shutil.move(str(source), final_path)
        elif mode == "reference":
            final_path = str(source)
            action = "referenced"

        logger.info(f"[INGEST] 文件操作完成: {final_path}")

        # 将文件直接交给 LightRAG 处理入库
        content_length = 0
        lightrag_result = None

        # 使用前面已检查的 KG 格式支持结果
        if not kg_ok:
            lr_status = "unsupported"
            lr_msg = kg_reason
            logger.info(f"[INGEST] KG unsupported: {kg_reason}")
            content_length = Path(final_path).stat().st_size
        else:
            lr_status = "skipped"
            lr_msg = ""
            try:
                from agent.tool_registry import get_registry

                registry = get_registry()
                insert_file_tool = registry.get("lightrag-server/lightrag_insert_file")
                if insert_file_tool:
                    lightrag_result = insert_file_tool(
                        file_path=str(Path(final_path).resolve()),
                        doc_id=str(Path(final_path).resolve()),
                    )
                    lr_status = lightrag_result.get("status", "unknown")
                    if lr_status == "ok":
                        logger.info(f"[INGEST] LightRAG insert_file: ok")
                    else:
                        lr_msg = lightrag_result.get("message", "")
                        logger.warning(
                            f"[INGEST] LightRAG insert_file failed: status={lr_status}, message={lr_msg}"
                        )
                    content_length = Path(final_path).stat().st_size
            except Exception as lr_err:
                logger.warning(f"[INGEST] LightRAG insert_file 失败（不影响文件入库）: {lr_err}")

        # Determine lightrag status for return value
        if lr_status == "unsupported":
            lr_final = "unsupported"
            lr_final_msg = lr_msg
        elif lightrag_result and lightrag_result.get("status") in ("success", "ok"):
            lr_final = "inserted"
            lr_final_msg = ""
        elif lightrag_result and lightrag_result.get("status") == "error":
            lr_final = "error"
            lr_final_msg = lightrag_result.get("message", "")
        else:
            lr_final = "skipped"
            lr_final_msg = ""

        # 根据知识图谱状态生成提示
        if lr_final == "unsupported":
            note = f"文件已存储，但{lr_final_msg}"
        elif lr_final == "error":
            note = f"文件已存储，知识图谱入库失败：{lr_final_msg}"
        else:
            note = "文件已存储，正在异步分析并存入知识图谱，处理时间可能较长，请耐心等待"

        return {
            "status": "success",
            "action": action,
            "file_path": str(Path(final_path).resolve()),
            "original_path": str(source),
            "category": category,
            "content_length": content_length,
            "lightrag": lr_final,
            "lightrag_message": lr_final_msg or None,
            "note": note,
            "kg_entities": [],
        }

    except PermissionError as e:
        logger.error(f"[INGEST] 权限错误: {e}")
        return {
            "status": "error",
            "error_code": "PERMISSION_DENIED",
            "message": f"无权限: {e}",
            "suggestion": "请检查文件或目录权限",
            "kg_entities": [],
        }
    except Exception as e:
        logger.exception(f"[INGEST] 失败: {e}")
        return {
            "status": "error",
            "error_code": "UNKNOWN_ERROR",
            "message": str(e),
            "suggestion": "请检查日志获取详细信息",
            "kg_entities": [],
        }


# ============== MCP Server ==============

server = Server("niu-photo-server")


@server.list_tools()
async def list_tools() -> list[Tool]:
    """列出可用工具"""
    return [
        Tool(
            name="ingest",
            description="""有状态统一入库工具 — 支持目录逐文件交互式入库

参数:
- path: 必填，文件路径或目录路径
- mode: copy（复制）| move（移动）| reference（引用），默认 copy
- category: 分类目录（文档需要分类时传入，照片不需要）
- action: start | interact | abort，默认空字符串

三阶段交互模式（目录入库）:
1. 初始化: ingest(path="E:/照片", action="start", mode="copy")
   → 扫描目录，创建会话，处理第一个文件，返回 progress/need_category
2. 中间态交互:
   - 继续（progress后）: ingest(path="E:/照片")
   - 回答分类（need_category后）: ingest(path="E:/照片", category="技术文档")
3. 中止: ingest(path="E:/照片", action="abort")

单文件入库（path是文件时）: 无状态，直接入库，action参数无效""",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "文件路径或目录路径"},
                    "mode": {
                        "type": "string",
                        "enum": ["copy", "move", "reference"],
                        "default": "copy",
                        "description": "文件操作模式",
                    },
                    "category": {"type": "string", "description": "文件分类目录。need_category状态时必须从available_categories中选择", "default": ""},
                    "action": {
                        "type": "string",
                        "enum": ["", "start", "interact", "abort"],
                        "default": "",
                        "description": "会话动作：start=初始化会话，interact/空=继续交互，abort=中止会话",
                    },
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="ingest_document",
            description="""文档入库工具 — 文件搬运 + 提交LightRAG异步处理

参数:
- file_path: 必填，源文件绝对路径
- category: 分类目录，不传则返回内容预览供判断分类
- mode: copy（复制）| move（移动）| reference（引用）

不传 category 时返回 need_category 状态+内容预览，判断分类后再次调用 ingest 或 ingest_document 传入 category。

返回:
- status: need_category | success | error
- action: created | versioned | renamed | referenced | skipped
- file_path: 存储后的完整路径
- lightrag: inserted | unsupported | skipped | error
- lightrag_message: unsupported/error 时的原因说明

冲突处理:
- 完全相同文件（哈希相同）→ 跳过
- 内容相似（语义相似度 > 阈值）→ 版本管理
- 内容不同 → 自动改名""",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "源文件绝对路径"},
                    "category": {
                        "type": "string",
                        "description": "文件分类目录（如：工作文档、个人资料、财务报告等）。不传则返回文件内容预览供你判断分类",
                        "default": "",
                    },
                    "mode": {
                        "type": "string",
                        "enum": ["copy", "move", "reference"],
                        "default": "copy",
                    },
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="ingest_photo",
            description="""照片入库工具（带人脸识别）

参数:
- file_path: 必填，照片文件绝对路径
- category: 分类（生活/工作/旅行/证件/其他），默认从 preferences.json 读取

返回:
- status: success | error
- photo_id: 照片唯一ID
- detected_persons: 检测到的人物列表 [{id, name, similarity}]
- abstract: L0 摘要（人物+时间）
- exif: EXIF 信息（taken_at, location, camera）

处理流程:
1. 提取 EXIF 信息（拍摄时间、GPS、相机）
2. 使用 InsightFace 检测人脸
3. 匹配已有人物或创建"未命名人物_N"
4. 生成 L0 摘要""",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "照片文件绝对路径"},
                    "category": {
                        "type": "string",
                        "description": "分类（生活/工作/旅行/证件/其他）",
                    },
                    "mode": {
                        "type": "string",
                        "enum": ["copy", "move", "reference"],
                        "default": "copy",
                        "description": "文件操作模式",
                    },
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="name_person",
            description="""为人物命名

参数:
- person_id: 必填，人物ID
- name: 必填，新名称

返回:
- status: success | error
- person_id: 人物ID
- name: 新名称
- auto_label: 自动标签（如"未命名人物_1"）""",
            inputSchema={
                "type": "object",
                "properties": {
                    "person_id": {"type": "string", "description": "人物ID"},
                    "name": {"type": "string", "description": "新名称"},
                },
                "required": ["person_id", "name"],
            },
        ),
        Tool(
            name="merge_persons",
            description="""合并两个人物

参数:
- person_a_id: 必填，保留的人物ID
- person_b_id: 必填，要合并到A的人物ID

返回:
- status: success | error
- merged_into: 合并到的人物ID
- name: 保留的名称
- photo_count: 合并后的照片数量
- deleted_person_id: 被删除的人物ID

说明:
- 保留 person_a 的名称
- 合并所有人脸向量，重新计算中心
- 更新所有照片关联
- 学习机制：如果相似度低于阈值，自动调整阈值""",
            inputSchema={
                "type": "object",
                "properties": {
                    "person_a_id": {"type": "string", "description": "保留的人物ID"},
                    "person_b_id": {"type": "string", "description": "要合并的人物ID"},
                },
                "required": ["person_a_id", "person_b_id"],
            },
        ),
        Tool(
            name="search_persons",
            description="""搜索人物（按名字模糊匹配）

参数:
- query: 搜索词（人名，子串匹配）
- limit: 返回数量（默认10）

返回:
- 匹配的人物列表""",
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "搜索词（人名）"},
                    "limit": {
                        "type": "integer",
                        "description": "返回数量",
                        "default": 10,
                    },
                },
                "required": ["query"],
            },
        ),
        Tool(
            name="get_unnamed_persons",
            description="""获取所有未命名人物

返回:
- 未命名人物列表，按出现次数排序
- 包含：id, auto_label, photo_count, has_valid_photos, photos: [{file_path, boxed_path}]
- boxed_path 是带人脸红框的图片路径，前端用 Markdown 图片语法 ![person_id|name](path) 显示
- has_valid_photos=false 表示该人物的照片文件已不存在""",
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),
        Tool(
            name="delete_person",
            description="""删除人物及其所有关联数据

警告：这会删除人物图谱中的节点，请谨慎使用。
只有在用户明确要求删除时才调用。

参数:
- person_id: 要删除的人物ID

返回:
- status: success | error
- message: 结果说明""",
            inputSchema={
                "type": "object",
                "properties": {
                    "person_id": {"type": "string", "description": "人物ID"},
                },
                "required": ["person_id"],
            },
        ),
        Tool(
            name="cleanup_deleted_photos",
            description="""清理已删除照片的数据库记录

在删除照片文件/目录后调用，清理数据库中的孤儿记录。
扫描 photos 表，删除文件不存在的记录及其关联的 faces 记录。

返回:
- deleted_photos: 删除的照片记录数
- deleted_faces: 删除的人脸记录数

使用场景：
- 用户删除了照片目录后
- 清理数据库中的残留记录""",
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),
        Tool(
            name="get_person_photos",
            description="""获取某人物的多张照片（用于"换一张"场景）

当用户说"看不清"、"换一张"时调用此工具。

参数:
- person_id: 人物ID
- limit: 最多返回几张照片（默认5）

返回:
- person_id, person_name
- photos: [{file_path, boxed_path, taken_at}, ...]

boxed_path 是带人脸红框的图片路径，前端用 Markdown 图片语法 ![person_id|name](path) 显示。""",
            inputSchema={
                "type": "object",
                "properties": {
                    "person_id": {"type": "string", "description": "人物ID"},
                    "limit": {
                        "type": "integer",
                        "description": "最多返回几张（默认5）",
                    },
                },
                "required": ["person_id"],
            },
        ),
        Tool(
            name="unload_face_model",
            description="""卸载人脸识别模型，释放内存（约 326MB）

在长时间空闲时调用此工具释放内存。
通常由系统在 SLEEP 状态时自动调用。

返回:
- status: success
- message: 卸载结果""",
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict[str, Any]) -> list[TextContent]:
    """调用工具 - 同步调用，避免 asyncio.to_thread 在 MCP stdio 环境中的问题"""
    try:
        logger.info(f"[CALL_TOOL] 开始: {name}")

        if name == "ingest":
            # 有状态统一入库工具
            result = ingest(
                path=arguments["path"],
                category=arguments.get("category", ""),
                mode=arguments.get("mode", "copy"),
                action=arguments.get("action", ""),
            )
        elif name == "ingest_document":
            result = ingest_document(
                file_path=arguments["file_path"],
                category=arguments.get("category", ""),
                mode=arguments.get("mode", "copy"),
            )
        elif name == "ingest_photo":
            result = ingest_photo(
                file_path=arguments["file_path"],
                category=arguments.get("category"),
                mode=arguments.get("mode", "copy"),
            )
        elif name == "name_person":
            result = name_person(
                person_id=arguments["person_id"],
                name=arguments["name"],
            )
        elif name == "merge_persons":
            result = merge_persons(
                person_a_id=arguments["person_a_id"],
                person_b_id=arguments["person_b_id"],
            )
        elif name == "delete_person":
            result = delete_person(person_id=arguments["person_id"])
        elif name == "cleanup_deleted_photos":
            result = cleanup_deleted_photos()
        elif name == "search_persons":
            persons = search_persons(
                query=arguments["query"],
                limit=arguments.get("limit", 10),
            )
            result = {
                "status": "success",
                "query": arguments["query"],
                "results": persons,
            }
        elif name == "get_unnamed_persons":
            persons = get_unnamed_persons()
            result = {"status": "success", "count": len(persons), "persons": persons}
        elif name == "get_person_photos":
            result = get_person_photos(
                person_id=arguments["person_id"],
                limit=arguments.get("limit", 5),
            )
        elif name == "unload_face_model":
            result = unload_face_model()
        else:
            result = {
                "status": "error",
                "error_code": "UNKNOWN_TOOL",
                "message": f"未知工具: {name}",
            }

        logger.info(f"[CALL_TOOL] 返回: {result.get('status')} {result.get('action')}")
        return [TextContent(type="text", text=json.dumps(result, ensure_ascii=False))]

    except Exception as e:
        logger.exception(f"[CALL_TOOL] 错误: {e}")
        return [
            TextContent(
                type="text",
                text=json.dumps(
                    {"status": "error", "message": str(e)}, ensure_ascii=False
                ),
            )
        ]


async def run_server():
    """运行 MCP 服务器"""
    # 使用共享向量服务，无需预加载 embedding model
    logger.info("[INIT] Photo server starting (using shared embedding service)")

    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream, write_stream, server.create_initialization_options()
        )


def preload_face_model():
    """预加载人脸识别模型（在 MCP stdio 启动前）"""
    import sys
    import os
    from pathlib import Path

    # 跨平台日志路径：使用 ~/.niu/logs
    log_dir = Path.home() / ".niu" / "logs"
    log_dir.mkdir(parents=True, exist_ok=True)
    log_file = log_dir / "photo-server-preload.log"

    # 强制输出到 stderr，确保在 MCP stdio 启动前可见
    with open(log_file, "a", encoding="utf-8") as f:
        f.write(f"[PRELOAD] Starting at {datetime.now().isoformat()}\n")
        f.flush()

    print("[PRELOAD] Pre-loading InsightFace and cv2...", file=sys.stderr, flush=True)
    logger.info("[PRELOAD] Pre-loading InsightFace and cv2...")

    try:
        print("[PRELOAD] Importing cv2...", file=sys.stderr, flush=True)
        import cv2

        print(
            f"[PRELOAD] cv2 imported, version: {cv2.__version__}",
            file=sys.stderr,
            flush=True,
        )

        print("[PRELOAD] Importing InsightFace...", file=sys.stderr, flush=True)
        from insightface.app import FaceAnalysis

        print("[PRELOAD] InsightFace imported", file=sys.stderr, flush=True)

        # 预加载模型（可选，看是否需要在启动时加载）
        # models_dir = get_models_dir()
        # face_model = FaceAnalysis(name="buffalo_l", root=str(models_dir), providers=["CPUExecutionProvider"])
        # face_model.prepare(ctx_id=-1)
        # global _face_model
        # _face_model = face_model

        print("[PRELOAD] Pre-load complete", file=sys.stderr, flush=True)
        logger.info("[PRELOAD] Pre-load complete")
    except Exception as e:
        print(f"[PRELOAD] Failed: {e}", file=sys.stderr, flush=True)
        logger.warning(f"[PRELOAD] Failed: {e}")


def main():
    """入口函数"""
    # 预加载 cv2 和 InsightFace 模块代码（不是模型）
    # 这是关键：在 MCP stdio 启动前预加载模块，避免动态导入时卡死
    # 注意：只导入模块，不加载模型（模型按需加载）
    preload_face_model()
    # 启动后台模型卸载定时器
    _start_model_unload_timer()
    asyncio.run(run_server())


if __name__ == "__main__":
    main()
